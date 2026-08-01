#!/usr/bin/env python3
"""Generate the PROTON Technical Proposal deck on the Devoteam 2025 brand template.

Content source: docs/proton-technical-proposal-2026-07-27.md (14-slide flow).
Styling mirrors the APL reference build (apl-azure-migration/build_proposal.py):
  - Montserrat everywhere; dark-grey body text, poppy reserved for accents.
  - Devoteam logo on every slide; photographic brand cover on page 1 ONLY.
  - Full brand palette (poppy / fire / aqua / beige / mint / blue) used with meaning.
The three architecture slides (6-8) are drawn NATIVELY in shapes (no embedded image).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml
from PIL import Image

BASE = "/Users/yudaadipratama/Archive/id-crm-ticketing/docs/proposals"
TEMPLATE = os.path.join(BASE, "templates/2025 Presentation template - Google Slides.pptx")
OUT = os.path.join(BASE, "PROTON - Technical Proposal - Self-Managed CRM on Google Cloud.pptx")
COVER = os.path.join(BASE, "assets/cover.jpg")
LOGO_DARK = os.path.join(BASE, "assets/logo_dark.png")    # poppy dot + dark wordmark (light bg)
LOGO_WHITE = os.path.join(BASE, "assets/logo_white.png")  # poppy circle + white wordmark (dark bg)

# ---- Brand palette (Devoteam 2025) ----
POPPY   = RGBColor(0xF8, 0x48, 0x5E)   # primary red
FIRE    = RGBColor(0xFC, 0xC3, 0x54)   # accent amber
PINK    = RGBColor(0xFC, 0xA2, 0xAE)   # secondary
POPPYLT = RGBColor(0xFD, 0xDA, 0xDE)   # poppy lighter
AQUA    = RGBColor(0xD7, 0xEB, 0xE7)   # aqua
BEIGE   = RGBColor(0xEF, 0xEA, 0xDC)   # beige
MINT    = RGBColor(0x5A, 0xB8, 0x91)   # fresh mint
DARK    = RGBColor(0x3C, 0x3C, 0x3A)   # dark grey (primary text)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xEF, 0xEE, 0xEE)   # light grey
MIDGREY = RGBColor(0x8A, 0x8A, 0x88)
BLUE    = RGBColor(0x4A, 0x8C, 0xCA)   # blue lagoon (brand)
BLUELT  = RGBColor(0xDA, 0xE8, 0xF4)   # soft blue tint
MINTLT  = RGBColor(0xDD, 0xEF, 0xE7)   # soft mint tint
FONT = "Montserrat"

prs = Presentation(TEMPLATE)
# strip all existing template slides (parts + rels), keep masters/layouts/theme
_pres_part = prs.part
_sldIdLst = prs.slides._sldIdLst
for sid in list(_sldIdLst):
    rId = sid.get(qn('r:id'))
    _sldIdLst.remove(sid)
    if rId:
        _pres_part.drop_rel(rId)

SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[0]  # no placeholders

_LOGO_RATIO = {}  # path -> h/w
def _logo_ratio(path):
    if path not in _LOGO_RATIO:
        w, h = Image.open(path).size
        _LOGO_RATIO[path] = h / w
    return _LOGO_RATIO[path]

# ---------- helpers ----------
def slide(bg=WHITE):
    """New slide. The template's base layout carries decorative shapes; slide shapes
    always render above layout shapes, so we lay down a solid canvas first."""
    sl = prs.slides.add_slide(BLANK)
    if bg is not None:
        sp = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        sp.shadow.inherit = False
        sp.fill.solid(); sp.fill.fore_color.rgb = bg
        sp.line.fill.background()
    return sl

def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    _set_fill(sp, color)
    if line is not None:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    return sp

def gradient_scrim(s, x, y, w, h, hexcolor="3C3C3A", ang_deg=0, a0=90, a1=0):
    """Directional transparent overlay (for photo legibility)."""
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.shadow.inherit = False
    spPr = sp._element.spPr
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    ang = int(ang_deg * 60000)
    grad = parse_xml(
        f'<a:gradFill {nsdecls("a")}><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{hexcolor}"><a:alpha val="{int(a0*1000)}"/></a:srgbClr></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{hexcolor}"><a:alpha val="{int(a1*1000)}"/></a:srgbClr></a:gs>'
        f'</a:gsLst><a:lin ang="{ang}" scaled="1"/></a:gradFill>')
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)
    sp.line.fill.background()
    return sp

def logo(s, dark=True, w=Inches(1.15), x=None, y=Inches(0.34)):
    path = LOGO_DARK if dark else LOGO_WHITE
    ww = int(w); hh = int(ww * _logo_ratio(path))
    if x is None:
        x = SW - Emu(ww) - Inches(0.5)
    s.shapes.add_picture(path, x, y, width=Emu(ww), height=Emu(hh))

def textbox(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=4, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, bold, color, *rest) in para:
            italic = rest[0] if rest else False
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT
            r.font.color.rgb = color; r.font.italic = italic
    return tb

def bullets(s, x, y, w, h, items, size=13, color=DARK, gap=6, marker="—", mcolor=POPPY,
            line_spacing=1.1):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0); p.line_spacing = line_spacing
        rm = p.add_run(); rm.text = f"{marker}  "
        rm.font.size = Pt(size); rm.font.bold = True; rm.font.name = FONT; rm.font.color.rgb = mcolor
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); r1.text = lead
            r1.font.size = Pt(size); r1.font.bold = True; r1.font.name = FONT; r1.font.color.rgb = color
            if rest:
                r2 = p.add_run(); r2.text = rest
                r2.font.size = Pt(size); r2.font.bold = False; r2.font.name = FONT; r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.bold = False; r.font.name = FONT; r.font.color.rgb = color
    return tb

def page_header(s, kicker, title):
    rect(s, Inches(0.55), Inches(0.42), Inches(0.28), Inches(0.14), POPPY)
    textbox(s, Inches(0.92), Inches(0.4), Inches(6.5), Inches(0.3),
            [[(kicker.upper(), 11, True, POPPY)]])
    tsize = 26 if len(title) <= 30 else (22 if len(title) <= 44 else 19)
    textbox(s, Inches(0.55), Inches(0.68), Inches(7.7), Inches(0.7),
            [[(title, tsize, True, DARK)]])
    rect(s, Inches(0.55), Inches(1.34), Inches(2.4), Pt(3), FIRE)
    logo(s, dark=True)

def footer(s, n):
    textbox(s, Inches(0.55), Inches(5.28), Inches(6.5), Inches(0.3),
            [[("Devoteam  ·  Technical Proposal  ·  PROTON — e.MAS Customer Operations", 8, False, MIDGREY)]])
    textbox(s, Inches(8.9), Inches(5.28), Inches(0.6), Inches(0.3),
            [[(str(n), 9, True, MIDGREY)]], align=PP_ALIGN.RIGHT)

def style_cell(cell, text, size=11, bold=False, color=DARK, fill=None,
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    cell.vertical_anchor = anchor
    cell.margin_left = Inches(0.09); cell.margin_right = Inches(0.09)
    cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for j, line in enumerate(text.split("\n")):
        pp = p if j == 0 else tf.add_paragraph()
        pp.alignment = align
        r = pp.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT; r.font.color.rgb = color

def make_table(s, x, y, w, rows, col_w, data, header_fill=DARK, header_color=WHITE,
               row_h=Inches(0.3), header_h=Inches(0.34), font=11, zebra=True,
               zebra_fill=BEIGE):
    ncols = len(col_w)
    nrows = len(data)
    gtbl = s.shapes.add_table(nrows, ncols, x, y, w, header_h + row_h*(nrows-1))
    tbl = gtbl.table
    tbl.first_row = False; tbl.horz_banding = False
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = cw
    tbl.rows[0].height = header_h
    for ri in range(1, nrows):
        tbl.rows[ri].height = row_h
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            if ri == 0:
                style_cell(cell, val, size=font, bold=True, color=header_color,
                           fill=header_fill, align=PP_ALIGN.LEFT)
            else:
                fill = zebra_fill if (zebra and ri % 2 == 0) else WHITE
                style_cell(cell, val, size=font, bold=False, color=DARK, fill=fill,
                           align=PP_ALIGN.LEFT)
    return tbl

def connector(s, x0, y0, x1, y1, color=POPPY, width=2.0, arrow=True):
    """Draw a straight H/V connector as a thin rectangle plus an optional triangle
    arrowhead. Uses only plain autoshapes — avoids the <p:cxnSp>/<a:tailEnd> XML
    that trips PowerPoint's repair prompt (python-pptx connectors)."""
    x0, y0, x1, y1 = int(x0), int(y0), int(x1), int(y1)
    lw = int(Pt(width)); ah = int(Inches(0.12))
    if y0 == y1:  # horizontal
        xs, xe = (x0, x1) if x1 >= x0 else (x1, x0)
        seg = max(0, (xe - xs) - (ah if arrow else 0))
        rect(s, Emu(xs), Emu(y0 - lw // 2), Emu(seg), Emu(lw), color)
        if arrow:
            t = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                   Emu(x1 - ah), Emu(y0 - ah // 2), Emu(ah), Emu(ah))
            t.rotation = 90 if x1 >= x0 else 270
            t.shadow.inherit = False; _set_fill(t, color)
    else:  # vertical
        ys, ye = (y0, y1) if y1 >= y0 else (y1, y0)
        seg = max(0, (ye - ys) - (ah if arrow else 0))
        top = ys if y1 >= y0 else ys + ah
        rect(s, Emu(x0 - lw // 2), Emu(top), Emu(lw), Emu(seg), color)
        if arrow:
            t = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                   Emu(x0 - ah // 2), Emu(y1 - ah // 2), Emu(ah), Emu(ah))
            t.rotation = 180 if y1 >= y0 else 0
            t.shadow.inherit = False; _set_fill(t, color)

def chip(s, x, y, w, h, title, sub, fill, tcolor=WHITE, scolor=None, tsize=11, ssize=8.5,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    """A labelled box for architecture diagrams."""
    rect(s, x, y, w, h, fill, shape=shape)
    if scolor is None:
        scolor = tcolor
    if sub:
        textbox(s, x + Inches(0.08), y + Inches(0.07), w - Inches(0.16), h - Inches(0.14),
                [[(title, tsize, True, tcolor)], [(sub, ssize, False, scolor)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=1)
    else:
        textbox(s, x + Inches(0.08), y + Inches(0.04), w - Inches(0.16), h - Inches(0.08),
                [[(title, tsize, True, tcolor)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

# ============================================================ SLIDES

# ---------- 1. COVER (only photographic slide) ----------
s = slide()
s.shapes.add_picture(COVER, 0, 0, width=SW, height=SH)
gradient_scrim(s, 0, 0, SW, SH, "3C3C3A", ang_deg=0, a0=92, a1=8)
rect(s, 0, 0, Inches(0.24), SH, POPPY)
textbox(s, Inches(0.7), Inches(0.85), Inches(8.6), Inches(0.4),
        [[("TECHNICAL PROPOSAL", 14, True, FIRE)]])
textbox(s, Inches(0.7), Inches(1.28), Inches(8.9), Inches(1.4),
        [[("PROTON Customer Complaint", 32, True, WHITE)],
         [("Management System", 32, True, WHITE)]],
        line_spacing=1.0, space_after=2)
textbox(s, Inches(0.7), Inches(2.74), Inches(8.7), Inches(0.5),
        [[("Self-managed, AI-native CRM on Google Cloud", 18, True, POPPY)]])
textbox(s, Inches(0.7), Inches(3.45), Inches(8.6), Inches(0.5),
        [[("Prepared for  ", 15, False, PINK), ("PROTON — e.MAS Customer Operations", 15, True, WHITE)]])
textbox(s, Inches(0.7), Inches(3.95), Inches(8.6), Inches(0.4),
        [[("Prepared by Devoteam  ·  built to your 5-channel complaint SOP", 12, False, WHITE)]])
textbox(s, Inches(0.7), Inches(4.38), Inches(8.6), Inches(0.4),
        [[("July 2026   ·   Version 1.0   ·   Confidential", 10, False, LIGHT)]])
logo(s, dark=False, w=Inches(1.6), x=Inches(0.68), y=SH - Inches(0.92))

# ---------- 2. EXECUTIVE SUMMARY ----------
s = slide()
page_header(s, "Section 1", "Executive Summary")
textbox(s, Inches(0.55), Inches(1.58), Inches(5.35), Inches(3.5),
        [[("The situation.  ", 13, True, POPPY),
          ("PROTON's complaint operation runs today on Zendesk — a per-seat SaaS CRM. Cost "
           "scales with every agent and premium feature, and the data and AI behaviour live inside "
           "a vendor you don't control.", 13, False, DARK)],
         [("The proposal.  ", 13, True, POPPY),
          ("Replace Zendesk with a self-managed CRM — Chatwoot Community (open-source, no per-seat "
           "licence) plus a Google Gemini / Vertex AI automation layer — deployed on Google Cloud "
           "and built directly to PROTON's own 5-channel process-flow SOP.", 13, False, DARK)]],
        line_spacing=1.2, space_after=10)
rect(s, Inches(6.2), Inches(1.58), Inches(3.25), Inches(3.55), AQUA)
rect(s, Inches(6.2), Inches(1.58), Inches(3.25), Inches(0.5), DARK)
textbox(s, Inches(6.4), Inches(1.66), Inches(3.0), Inches(0.4),
        [[("WHY NOW / WHY US", 12, True, WHITE)]])
bullets(s, Inches(6.4), Inches(2.28), Inches(2.9), Inches(2.8),
        [("Cost — ", "eliminate per-agent SaaS licensing; pay for infrastructure, not seats"),
         ("Ownership — ", "your conversations, KB, AI prompts & data, inside PROTON's GCP tenancy"),
         ("Speed — ", "majority already built & demoed live against your SOP — migration, not greenfield"),
         ("AI-native — ", "Gemini answers in the customer's language, grounded on your FAQ/KB")],
        size=11, gap=9, mcolor=POPPY)
footer(s, 2)

# ---------- 3. OUR UNDERSTANDING ----------
s = slide()
page_header(s, "Section 1", "Our Understanding of Your Requirements")
textbox(s, Inches(0.55), Inches(1.5), Inches(8.9), Inches(0.3),
        [[("Mapped end-to-end from your Customer Complaint Management System requirement and the "
           "CRM Process Flow workbook (WhatsApp / Social / Email / IVR / SSI).", 10, False, MIDGREY, True)]])
rows = [["#", "Requirement cluster", "What you need"],
        ["1", "Omnichannel inbound", "Call / Email / WhatsApp / Social in one agent view, single customer record, new-message alerts, voice-to-text"],
        ["2", "Agent management", "On-duty check before escalation, per-agent channel priorities, status-aware auto-assign, reminders & timeouts"],
        ["3", "FAQ + AI support", "Live-editable knowledge base, AI auto-suggested replies from FAQ, one-click reference, FAQ quality scoring"],
        ["4", "Escalation & SLA", "Rules → PIC by category/SOP, email + WhatsApp alerts, case states (WIP/Resolved/Temp-Closed), 8h/48h auto-escalation"],
        ["5", "Customer 360", "Two-way DMS + TSP integration, auto-identify by number, 360 view card (personal / vehicle / service / call history)"],
        ["6", "Reporting / BI", "Channel & division analytics, dept/PIC & Call-Centre KPIs, NPS, lifecycle, scheduled exports, anomaly alerts, Power BI"]]
make_table(s, Inches(0.55), Inches(1.8), Inches(8.9), len(rows),
           [Inches(0.4), Inches(2.1), Inches(6.4)], rows,
           row_h=Inches(0.44), header_h=Inches(0.3), font=9, zebra_fill=AQUA)
textbox(s, Inches(0.55), Inches(4.86), Inches(8.9), Inches(0.3),
        [[("Cross-cutting:  ", 9, True, POPPY),
          ("reduce SaaS cost · own the data & AI · honour the 5-channel SOP timers and surveys exactly · multi-tenant.",
           9, False, DARK)]])
footer(s, 3)

# ---------- 4. WHAT TO EXPECT ----------
s = slide()
page_header(s, "Section 1", "What to Expect")
rows = [["Business Expectation", "Technical Expectation"],
        ["No per-agent / per-seat licence fees", "Chatwoot Community (open-source), self-hosted on GCP — unlimited agents"],
        ["Your data stays yours", "All conversations, KB and AI prompts inside PROTON's GCP project"],
        ["AI answers customers instantly, in their language", "Gemini on Vertex AI, same-language replies, grounded on your FAQ/KB"],
        ["Your SOP, not a generic product", "Process flow as code: disclaimers, idle-close, YES/NO resolution, rating surveys"],
        ["Faster complaint resolution", "Auto-classification → PIC routing, SLA timers (2-min WA ack … 48h alarm)"],
        ["Management visibility", "BigQuery analytics + Looker / Power BI, scheduled PDF/Excel, anomaly alerts"],
        ["Predictable, scalable cost", "Managed GCP services; pay for compute/storage consumed, scale per tenant"],
        ["Low delivery risk & continuity", "Majority already built & demoed live; phased cutover, parallel-run first"]]
make_table(s, Inches(0.55), Inches(1.55), Inches(8.9), len(rows),
           [Inches(3.9), Inches(5.0)], rows,
           row_h=Inches(0.4), header_h=Inches(0.34), font=10.5)
footer(s, 4)

# ---------- 5. SECTION 2 DIVIDER ----------
_DIV_TINT = {1: AQUA, 2: BEIGE, 3: POPPYLT, 4: BLUELT}
def divider(num, title, sub):
    s = slide()
    rect(s, 0, 0, SW, SH, WHITE)
    textbox(s, Inches(5.1), Inches(1.15), Inches(5.2), Inches(3.4),
            [[(f"0{num}", 200, True, _DIV_TINT[num])]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.7), Inches(1.72), Inches(0.28), Inches(0.14), POPPY)
    textbox(s, Inches(1.07), Inches(1.7), Inches(5.0), Inches(0.3),
            [[(f"SECTION 0{num}", 12, True, POPPY)]])
    textbox(s, Inches(0.7), Inches(2.05), Inches(6.4), Inches(1.4),
            [[(title, 36, True, DARK)]], line_spacing=1.0)
    rect(s, Inches(0.72), Inches(3.62), Inches(2.6), Pt(3), FIRE)
    textbox(s, Inches(0.72), Inches(3.84), Inches(6.4), Inches(0.6),
            [[(sub, 14, False, MIDGREY)]])
    logo(s, dark=True)
    return s

divider(2, "Target Architecture", "Zendesk (SaaS)  →  self-managed on Google Cloud")

# ---------- 6. CURRENT STATE vs TARGET STATE (native before/after) ----------
s = slide()
page_header(s, "Section 2", "Current State vs. Target State")

# BEFORE panel
bx = Inches(0.55); bw = Inches(3.55); by = Inches(1.72); bh = Inches(3.05)
rect(s, bx, by, bw, bh, LIGHT)
rect(s, bx, by, bw, Inches(0.44), MIDGREY)
textbox(s, bx, by + Inches(0.06), bw, Inches(0.34),
        [[("TODAY — ZENDESK (SaaS)", 12, True, WHITE)]], align=PP_ALIGN.CENTER)
bullets(s, bx + Inches(0.22), by + Inches(0.6), bw - Inches(0.44), Inches(2.4),
        ["Per-agent subscription; SLA / roles / AI behind paid tiers",
         "Data & AI behaviour hosted by the vendor",
         "FAQ / Guide inside Zendesk — limited AI-grounding control",
         "DMS/TSP & telephony bolted onto a closed platform"],
        size=10.5, gap=9, marker="✕", mcolor=MIDGREY)

# migrate arrow (center)
midx = Inches(4.35)
rect(s, midx, Inches(2.86), Inches(0.7), Inches(0.78), POPPY, shape=MSO_SHAPE.RIGHT_ARROW)
textbox(s, Inches(4.2), Inches(2.5), Inches(1.0), Inches(0.3),
        [[("MIGRATE", 9.5, True, POPPY)]], align=PP_ALIGN.CENTER)

# AFTER panel
ax = Inches(5.35); aw = Inches(4.1)
rect(s, ax, by, aw, bh, AQUA)
rect(s, ax, by, aw, Inches(0.44), DARK)
textbox(s, ax, by + Inches(0.06), aw, Inches(0.34),
        [[("TARGET — SELF-MANAGED ON GCP", 12, True, WHITE)]], align=PP_ALIGN.CENTER)
bullets(s, ax + Inches(0.22), by + Inches(0.6), aw - Inches(0.44), Inches(2.4),
        ["Open-source Chatwoot core — no seat licences; SLA / roles / AI built, not rented",
         "Runs inside PROTON's GCP project — full data residency & control",
         "KB + AI prompts owned and edited by PROTON operators (no-code)",
         "Native first-party integration surfaces for DMS/TSP, telephony & BI"],
        size=10.5, gap=9, marker="✓", mcolor=MINT)
textbox(s, Inches(0.55), Inches(4.95), Inches(8.9), Inches(0.35),
        [[("Same capabilities — moved onto foundations PROTON owns and controls, at infrastructure cost.",
           10.5, True, DARK)]], align=PP_ALIGN.CENTER)
footer(s, 6)

# ---------- 7. GCP PRODUCTION ARCHITECTURE (native stack) ----------
s = slide()
page_header(s, "Section 2", "GCP Production Architecture")

# Internet -> LB
textbox(s, Inches(0.55), Inches(1.62), Inches(1.0), Inches(0.4),
        [[("Internet", 10, True, MIDGREY)]], anchor=MSO_ANCHOR.MIDDLE)
lb_x, lb_y, lb_w, lb_h = Inches(1.55), Inches(1.6), Inches(1.65), Inches(0.5)
chip(s, lb_x, lb_y, lb_w, lb_h, "Cloud LB + Caddy", "TLS · routing · per-tenant entry", POPPY, tsize=10.5, ssize=8)
connector(s, Inches(1.42), Inches(1.85), lb_x, Inches(1.85), color=MIDGREY)

# Application tier (3 services)
app_y = Inches(2.45); app_h = Inches(0.62); app_w = Inches(2.7)
apps = [("Chatwoot  (GKE)", "Rails + Sidekiq · CRM / live-chat core", Inches(0.75)),
        ("agent  (Cloud Run)", "Webhook sync · AI orchestration · SLA", Inches(3.6)),
        ("backend  (Cloud Run)", "Gemini agent · KB · routing · metrics", Inches(6.45))]
rect(s, Inches(0.55), Inches(2.3), Inches(8.9), Inches(0.94), MINTLT)
textbox(s, Inches(0.62), Inches(2.32), Inches(3.0), Inches(0.2),
        [[("APPLICATION  (containerized, autoscaling)", 8, True, MINT)]])
for t, d, x in apps:
    chip(s, x, app_y, app_w, app_h, t, d, MINT, tsize=10.5, ssize=7.5)
    connector(s, x + app_w/2, lb_y + lb_h, x + app_w/2, app_y, color=MIDGREY, arrow=False, width=1.25)

# Data + AI tier
data_y = Inches(3.6); data_h = Inches(0.66); dw = Inches(2.12)
rect(s, Inches(0.55), Inches(3.46), Inches(8.9), Inches(0.98), BLUELT)
textbox(s, Inches(0.62), Inches(3.48), Inches(4.0), Inches(0.2),
        [[("AI · DATA · STATE  (managed GCP services)", 8, True, BLUE)]])
tiles = [("Vertex AI — Gemini", "drafts · classify · same-language", BLUE),
         ("Vertex Search + pgvector", "KB grounding — no black box", BLUE),
         ("Cloud SQL (Postgres)", "Chatwoot DB + per-tenant KB", DARK),
         ("Memorystore · Storage", "Redis queues · attachments", DARK)]
tx = Inches(0.62)
for t, d, c in tiles:
    chip(s, tx, data_y, dw, data_h, t, d, c, tsize=9.5, ssize=7.5)
    tx = Emu(int(tx) + int(dw) + int(Inches(0.13)))

# Bottom strip: analytics + ops
rect(s, Inches(0.55), Inches(4.62), Inches(8.9), Inches(0.5), BEIGE)
textbox(s, Inches(0.72), Inches(4.62), Inches(8.6), Inches(0.5),
        [[("Secret Manager  ·  BigQuery → Looker Studio / Power BI  ·  Cloud Monitoring & Logging  ·  "
           "Automated backups  ·  Multi-tenant isolation (one app stack + isolated DBs per business unit)",
           9, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 7)

# ---------- 8. DATA MIGRATION (native flow) ----------
s = slide()
page_header(s, "Section 2", "Data Migration — Zendesk → Platform")

# top flow: export -> transform -> import
flow = [("Zendesk API\nexport", MIDGREY), ("Transform &\nmap", FIRE),
        ("Bulk import\nChatwoot API", MINT), ("KB ingest\nVertex + pgvector", BLUE)]
fw = Inches(1.9); fh = Inches(0.8); fy = Inches(1.85); gap = Inches(0.32)
fx = Inches(0.55)
for i, (t, c) in enumerate(flow):
    chip(s, fx, fy, fw, fh, t.split("\n")[0], t.split("\n")[1], c,
         tcolor=(DARK if c is FIRE else WHITE), scolor=(DARK if c is FIRE else WHITE),
         tsize=11, ssize=9)
    if i < len(flow) - 1:
        ax0 = Emu(int(fx) + int(fw))
        connector(s, ax0, fy + fh/2, Emu(int(ax0) + int(gap)), fy + fh/2, color=POPPY)
    fx = Emu(int(fx) + int(fw) + int(gap))

# what we migrate (left) / how (right)
textbox(s, Inches(0.55), Inches(3.05), Inches(4.4), Inches(0.3), [[("WHAT WE MIGRATE", 12, True, POPPY)]])
bullets(s, Inches(0.55), Inches(3.45), Inches(4.4), Inches(1.7),
        [("Tickets / conversations ", "→ Chatwoot conversations (status & labels)"),
         ("Contacts / customers ", "→ Chatwoot contacts (phone/email, custom fields)"),
         ("Knowledge / Guide (FAQ) ", "→ PROTON-owned KB (Vertex Search + pgvector)"),
         ("Attachments → Storage ", "· reporting history → BigQuery")],
        size=10.5, gap=7)
textbox(s, Inches(5.15), Inches(3.05), Inches(4.3), Inches(0.3), [[("HOW — DE-RISKED, NO BIG-BANG", 12, True, POPPY)]])
bullets(s, Inches(5.15), Inches(3.45), Inches(4.3), Inches(1.7),
        [("Export → transform ", "→ bulk import via Chatwoot API + KB ingest"),
         ("Parallel run ", "— platform live alongside Zendesk during validation"),
         ("Cutover ", "— switch channel endpoints once data fidelity validated"),
         ("Decommission ", "Zendesk after sign-off")],
        size=10.5, gap=7, mcolor=MINT)
footer(s, 8)

# ---------- 9. SECTION 3 DIVIDER ----------
divider(3, "Delivery", "Timeline · Scope · Deliverables · Why this works")

# ---------- 10. PROJECT TIMELINE ----------
s = slide()
page_header(s, "Section 3", "Project Timeline — ~12–16 Weeks, Phased")
rows = [["Phase", "Weeks", "Focus", "Key outcomes"],
        ["P0 — Discovery &\nGCP foundation", "1–2", "Access, data contracts, GCP project / landing zone",
         "GCP org/project, network, Secret Manager, CI; Zendesk export sample validated"],
        ["P1 — Core platform\ncutover", "3–7", "Stand up production stack; migrate data; wire channels",
         "Chatwoot + agent + backend live on GKE/Cloud Run + Cloud SQL; WA/Email/Social connected; KB migrated; AI + SOP flows live"],
        ["P2 — Gap closure", "6–12", "The genuine net-new work (parallel with P1 tail)",
         "Customer 360 + DMS/TSP; RBAC roles/permissions; telephony/IVR hookup; Reports Tier-2 tabs"],
        ["P3 — Hardening,\nparallel-run & handover", "12–16", "HA, monitoring, UAT, cutover, decommission",
         "Autoscaling/HA, backups, dashboards; UAT sign-off; Zendesk decommissioned; ops runbook + training"]]
make_table(s, Inches(0.55), Inches(1.6), Inches(8.9), len(rows),
           [Inches(1.75), Inches(0.7), Inches(2.85), Inches(3.6)], rows,
           row_h=Inches(0.74), header_h=Inches(0.32), font=9.5, zebra_fill=BLUELT)
textbox(s, Inches(0.55), Inches(5.02), Inches(8.9), Inches(0.3),
        [[("P0/P1 largely assemble already-built components; real engineering concentrates in P2 (DMS/TSP, telephony — needs PROTON access).",
           9, False, MIDGREY, True)]])
footer(s, 10)

# ---------- 11. SCOPE OF WORK ----------
s = slide()
page_header(s, "Section 3", "Scope of Work")
sc1 = ["GCP landing zone: project, networking, IAM, Secret Manager, CI/CD",
       "Deploy production Chatwoot + agent + backend on GKE / Cloud Run + Cloud SQL",
       "Omnichannel: WhatsApp, Email, Social (FB/IG), Web widget, Voice-to-text",
       "AI layer: Gemini/Vertex reply drafting, same-language answers, KB grounding, auto-classification",
       "Knowledge base: migration + operator no-code authoring (FAQ + pgvector)",
       "Agent management: presence, channel-priority routing, status-aware assignment, My-Tasks timers"]
sc2 = ["Escalation & SLA: category→PIC routing, email + WhatsApp alerts, case states, 8h/48h auto-escalation",
       "Lifecycle/SOP: disclaimers, idle-close, YES/NO gate, rating surveys, auto-categorization, email auto-ack",
       "Customer 360: DMS + TSP two-way integration + 360 view card (subject to API access)",
       "RBAC: per-tenant roles & permissions model",
       "Reporting/BI: BigQuery warehouse + views, Looker/Power BI, scheduled exports, anomaly alerts, NPS/CSAT",
       "Data migration from Zendesk + Ops: monitoring, backups, multi-tenant isolation, runbook + training"]
bullets(s, Inches(0.55), Inches(1.6), Inches(4.45), Inches(3.5), sc1, size=10.5, gap=9)
bullets(s, Inches(5.15), Inches(1.6), Inches(4.3), Inches(3.5), sc2, size=10.5, gap=9)
textbox(s, Inches(0.55), Inches(5.02), Inches(8.9), Inches(0.28),
        [[("Assumptions:  ", 9, True, POPPY),
          ("PROTON provides GCP billing/org access, DMS/TSP API access, Zendesk export, and a telephony/CTI provider decision.",
           9, False, DARK)]])
footer(s, 11)

# ---------- 12. OUT OF SCOPE ----------
s = slide()
page_header(s, "Section 3", "Out of Scope")
oo1 = ["Telephony/PSTN carrier procurement — we integrate the CTI/provider; the phone-line contract & number provisioning is PROTON's (voice-to-text engine itself is built)",
       "DMS / TSP source-system changes — we consume their APIs; internal changes are the owners' responsibility",
       "SSI dealer-satisfaction survey workflow — lives in the e.MAS app; the platform ingests & reports SSI, but does not host that flow"]
oo2 = ["Power BI licences — we connect to BigQuery; Microsoft licensing is PROTON's",
       "Non-CRM business applications, custom hardware, and end-user device management",
       "Ongoing FAQ/KB content authoring after handover — enabled via no-code tools, owned by PROTON operators",
       "Third-party SaaS subscriptions beyond the GCP + open-source stack"]
bullets(s, Inches(0.55), Inches(1.7), Inches(4.45), Inches(3.4), oo1, size=11, gap=16, marker="✕", mcolor=MIDGREY)
bullets(s, Inches(5.15), Inches(1.7), Inches(4.3), Inches(3.4), oo2, size=11, gap=16, marker="✕", mcolor=MIDGREY)
footer(s, 12)

# ---------- 13. DELIVERABLES ----------
s = slide()
page_header(s, "Section 3", "Deliverables")
rows = [["Deliverable", "Description"],
        ["Production GCP environment", "Fully provisioned, IaC-described GCP project running the platform"],
        ["Migrated CRM", "Chatwoot with imported Zendesk tickets, contacts, KB, attachments"],
        ["AI automation layer", "Gemini/Vertex reply, classification, KB-grounded answers, SOP flows — live"],
        ["Customer 360 integration", "DMS/TSP two-way sync + 360 view card (pending PROTON API access)"],
        ["RBAC", "Roles & permissions model with admin UI"],
        ["Reporting suite", "BigQuery + Looker/Power BI dashboards, scheduled exports, anomaly alerts"],
        ["Data-migration report", "Record counts, fidelity validation, parallel-run results"],
        ["Operations runbook", "Deploy, backup/restore, monitoring, tenant provisioning, incident response"],
        ["Admin & agent training", "Sessions + materials for operators and administrators"],
        ["UAT sign-off & cutover", "Validated go-live and Zendesk decommission"]]
make_table(s, Inches(0.55), Inches(1.55), Inches(8.9), len(rows),
           [Inches(3.1), Inches(5.8)], rows,
           row_h=Inches(0.335), header_h=Inches(0.32), font=10.5, zebra_fill=MINTLT)
footer(s, 13)

# ---------- 14. WHY THIS WORKS (dark closing) ----------
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, 0, Inches(0.24), SH, POPPY)
rect(s, 0, Inches(4.86), SW, Inches(0.10), FIRE)
textbox(s, Inches(0.7), Inches(0.72), Inches(8.6), Inches(0.9),
        [[("Why This Works", 34, True, WHITE)]])
rect(s, Inches(0.72), Inches(1.48), Inches(2.4), Pt(3), FIRE)
bullets(s, Inches(0.75), Inches(1.82), Inches(8.5), Inches(2.6),
        [("Proven — ", "the platform already runs live against PROTON's SOP: migration + hardening, not invention"),
         ("Open + owned — ", "no per-seat lock-in; PROTON controls the data, the AI, and the roadmap"),
         ("AI-native on Google — ", "Gemini + Vertex AI, grounded on PROTON's own knowledge; measurable deflection & faster resolution"),
         ("Cost-down, capability-up — ", "replace recurring SaaS licensing with managed infrastructure you scale on your terms")],
        size=13, gap=11, color=WHITE, mcolor=FIRE)
textbox(s, Inches(0.7), Inches(4.5), Inches(8.6), Inches(0.32),
        [[("Cost  ·  Ownership  ·  Speed", 14, True, PINK)]])
logo(s, dark=False, w=Inches(1.5), x=SW - Inches(2.0), y=SH - Inches(0.9))

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
