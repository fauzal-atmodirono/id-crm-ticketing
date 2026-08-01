#!/usr/bin/env python3
"""PROTON Technical Proposal deck — v2 (expanded, requirement-grounded).

Sources folded in:
  - proton-technical-proposal-2026-07-27.md            (baseline narrative)
  - CRM System Enhancement 260414.pdf  (PROTON BRD)    (6 goals, 360, reporting)
  - CRM Process Flow (1).xlsx          (SOP workbook)   (5-channel timers/SLA)
  - Proton x Devoteam CRM System Update transcript      (28 Jul demo + gaps + commercials)

Style mirrors the APL reference build (Devoteam 2025 brand). Native in-slide
diagrams (no embedded images). Output is a NEW file; v1 is left untouched.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml
from PIL import Image

BASE = "/Users/yudaadipratama/Archive/id-crm-ticketing/docs/proposals"
TEMPLATE = os.path.join(BASE, "templates/2025 Presentation template - Google Slides.pptx")
OUT = os.path.join(BASE, "PROTON - Technical Proposal - Self-Managed CRM on Google Cloud (v2).pptx")
COVER = os.path.join(BASE, "assets/cover.jpg")
LOGO_DARK = os.path.join(BASE, "assets/logo_dark.png")
LOGO_WHITE = os.path.join(BASE, "assets/logo_white.png")

# ---- Brand palette (Devoteam 2025) ----
POPPY   = RGBColor(0xF8, 0x48, 0x5E)
FIRE    = RGBColor(0xFC, 0xC3, 0x54)
PINK    = RGBColor(0xFC, 0xA2, 0xAE)
POPPYLT = RGBColor(0xFD, 0xDA, 0xDE)
AQUA    = RGBColor(0xD7, 0xEB, 0xE7)
BEIGE   = RGBColor(0xEF, 0xEA, 0xDC)
MINT    = RGBColor(0x5A, 0xB8, 0x91)
DARK    = RGBColor(0x3C, 0x3C, 0x3A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xEF, 0xEE, 0xEE)
MIDGREY = RGBColor(0x8A, 0x8A, 0x88)
BLUE    = RGBColor(0x4A, 0x8C, 0xCA)
BLUELT  = RGBColor(0xDA, 0xE8, 0xF4)
MINTLT  = RGBColor(0xDD, 0xEF, 0xE7)
FONT = "Montserrat"

prs = Presentation(TEMPLATE)
_pres_part = prs.part
_sldIdLst = prs.slides._sldIdLst
for sid in list(_sldIdLst):
    rId = sid.get(qn('r:id'))
    _sldIdLst.remove(sid)
    if rId:
        _pres_part.drop_rel(rId)

SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[0]

_LOGO_RATIO = {}
def _logo_ratio(path):
    if path not in _LOGO_RATIO:
        w, h = Image.open(path).size
        _LOGO_RATIO[path] = h / w
    return _LOGO_RATIO[path]

# ---------- helpers ----------
def slide(bg=WHITE):
    sl = prs.slides.add_slide(BLANK)
    if bg is not None:
        sp = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        sp.shadow.inherit = False
        sp.fill.solid(); sp.fill.fore_color.rgb = bg
        sp.line.fill.background()
    return sl

def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    _set_fill(sp, color)
    if line is not None:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    return sp

def gradient_scrim(s, x, y, w, h, hexcolor="3C3C3A", ang_deg=0, a0=90, a1=0):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.shadow.inherit = False
    spPr = sp._element.spPr
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    ang = int(ang_deg * 60000)
    grad = parse_xml(
        f'<a:gradFill {nsdecls("a")}><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{hexcolor}"><a:alpha val="{int(a0*1000)}"/></a:srgbClr></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{hexcolor}"><a:alpha val="{int(a1*1000)}"/></a:srgbClr></a:gs>'
        f'</a:gsLst><a:lin ang="{ang}" scaled="1"/></a:gradFill>')
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)
    sp.line.fill.background()
    return sp

def logo(s, dark=True, w=Inches(1.15), x=None, y=Inches(0.34)):
    path = LOGO_DARK if dark else LOGO_WHITE
    ww = int(w); hh = int(ww * _logo_ratio(path))
    if x is None:
        x = SW - Emu(ww) - Inches(0.5)
    s.shapes.add_picture(path, x, y, width=Emu(ww), height=Emu(hh))

def textbox(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=4, line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, bold, color, *rest) in para:
            italic = rest[0] if rest else False
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT
            r.font.color.rgb = color; r.font.italic = italic
    return tb

def bullets(s, x, y, w, h, items, size=13, color=DARK, gap=6, marker="—", mcolor=POPPY,
            line_spacing=1.1):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0); p.line_spacing = line_spacing
        rm = p.add_run(); rm.text = f"{marker}  "
        rm.font.size = Pt(size); rm.font.bold = True; rm.font.name = FONT; rm.font.color.rgb = mcolor
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); r1.text = lead
            r1.font.size = Pt(size); r1.font.bold = True; r1.font.name = FONT; r1.font.color.rgb = color
            if rest:
                r2 = p.add_run(); r2.text = rest
                r2.font.size = Pt(size); r2.font.bold = False; r2.font.name = FONT; r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.bold = False; r.font.name = FONT; r.font.color.rgb = color
    return tb

def page_header(s, kicker, title):
    rect(s, Inches(0.55), Inches(0.42), Inches(0.28), Inches(0.14), POPPY)
    textbox(s, Inches(0.92), Inches(0.4), Inches(6.5), Inches(0.3),
            [[(kicker.upper(), 11, True, POPPY)]])
    tsize = 26 if len(title) <= 30 else (22 if len(title) <= 44 else 19)
    textbox(s, Inches(0.55), Inches(0.68), Inches(7.9), Inches(0.7),
            [[(title, tsize, True, DARK)]])
    rect(s, Inches(0.55), Inches(1.34), Inches(2.4), Pt(3), FIRE)
    logo(s, dark=True)

def footer(s, n):
    textbox(s, Inches(0.55), Inches(5.28), Inches(6.7), Inches(0.3),
            [[("Devoteam  ·  Technical Proposal  ·  PROTON — e.MAS Customer Operations", 8, False, MIDGREY)]])
    textbox(s, Inches(8.9), Inches(5.28), Inches(0.6), Inches(0.3),
            [[(str(n), 9, True, MIDGREY)]], align=PP_ALIGN.RIGHT)

def style_cell(cell, text, size=11, bold=False, color=DARK, fill=None,
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    cell.vertical_anchor = anchor
    cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for j, line in enumerate(text.split("\n")):
        pp = p if j == 0 else tf.add_paragraph()
        pp.alignment = align
        r = pp.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT; r.font.color.rgb = color

def make_table(s, x, y, w, rows, col_w, data, header_fill=DARK, header_color=WHITE,
               row_h=Inches(0.3), header_h=Inches(0.34), font=11, zebra=True,
               zebra_fill=BEIGE):
    ncols = len(col_w); nrows = len(data)
    gtbl = s.shapes.add_table(nrows, ncols, x, y, w, header_h + row_h*(nrows-1))
    tbl = gtbl.table
    tbl.first_row = False; tbl.horz_banding = False
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = cw
    tbl.rows[0].height = header_h
    for ri in range(1, nrows):
        tbl.rows[ri].height = row_h
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            if ri == 0:
                style_cell(cell, val, size=font, bold=True, color=header_color,
                           fill=header_fill, align=PP_ALIGN.LEFT)
            else:
                fill = zebra_fill if (zebra and ri % 2 == 0) else WHITE
                style_cell(cell, val, size=font, bold=False, color=DARK, fill=fill,
                           align=PP_ALIGN.LEFT)
    return tbl

def connector(s, x0, y0, x1, y1, color=POPPY, width=2.0, arrow=True):
    """H/V connector as a thin rect + optional triangle arrowhead (no cxnSp)."""
    x0, y0, x1, y1 = int(x0), int(y0), int(x1), int(y1)
    lw = int(Pt(width)); ah = int(Inches(0.12))
    if y0 == y1:
        xs, xe = (x0, x1) if x1 >= x0 else (x1, x0)
        seg = max(0, (xe - xs) - (ah if arrow else 0))
        rect(s, Emu(xs), Emu(y0 - lw // 2), Emu(seg), Emu(lw), color)
        if arrow:
            t = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                   Emu(x1 - ah), Emu(y0 - ah // 2), Emu(ah), Emu(ah))
            t.rotation = 90 if x1 >= x0 else 270
            t.shadow.inherit = False; _set_fill(t, color)
    else:
        ys, ye = (y0, y1) if y1 >= y0 else (y1, y0)
        seg = max(0, (ye - ys) - (ah if arrow else 0))
        top = ys if y1 >= y0 else ys + ah
        rect(s, Emu(x0 - lw // 2), Emu(top), Emu(lw), Emu(seg), color)
        if arrow:
            t = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                   Emu(x0 - ah // 2), Emu(y1 - ah // 2), Emu(ah), Emu(ah))
            t.rotation = 180 if y1 >= y0 else 0
            t.shadow.inherit = False; _set_fill(t, color)

def chip(s, x, y, w, h, title, sub, fill, tcolor=WHITE, scolor=None, tsize=11, ssize=8.5,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    rect(s, x, y, w, h, fill, shape=shape)
    if scolor is None:
        scolor = tcolor
    if sub:
        textbox(s, x + Inches(0.06), y + Inches(0.06), w - Inches(0.12), h - Inches(0.12),
                [[(title, tsize, True, tcolor)], [(sub, ssize, False, scolor)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=1)
    else:
        textbox(s, x + Inches(0.06), y + Inches(0.04), w - Inches(0.12), h - Inches(0.08),
                [[(title, tsize, True, tcolor)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

_PG = [0]
def npg():
    _PG[0] += 1
    return _PG[0]

_DIV_TINT = {2: BEIGE, 3: POPPYLT, 4: BLUELT, 5: MINTLT}
def divider(num, title, sub):
    s = slide()
    rect(s, 0, 0, SW, SH, WHITE)
    textbox(s, Inches(5.1), Inches(1.15), Inches(5.2), Inches(3.4),
            [[(f"0{num}", 200, True, _DIV_TINT[num])]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.7), Inches(1.72), Inches(0.28), Inches(0.14), POPPY)
    textbox(s, Inches(1.07), Inches(1.7), Inches(5.0), Inches(0.3),
            [[(f"SECTION 0{num}", 12, True, POPPY)]])
    tsize = 38 if len(title) <= 26 else (30 if len(title) <= 40 else 25)
    textbox(s, Inches(0.7), Inches(2.05), Inches(8.3), Inches(1.5),
            [[(title, tsize, True, DARK)]], line_spacing=1.0)
    rect(s, Inches(0.72), Inches(3.62), Inches(2.6), Pt(3), FIRE)
    textbox(s, Inches(0.72), Inches(3.84), Inches(8.3), Inches(0.6),
            [[(sub, 14, False, MIDGREY)]])
    logo(s, dark=True)
    npg()
    return s

# ============================================================ SLIDES

# ---------- 1. COVER ----------
s = slide()
s.shapes.add_picture(COVER, 0, 0, width=SW, height=SH)
gradient_scrim(s, 0, 0, SW, SH, "3C3C3A", ang_deg=0, a0=92, a1=8)
rect(s, 0, 0, Inches(0.24), SH, POPPY)
textbox(s, Inches(0.7), Inches(0.82), Inches(8.6), Inches(0.4),
        [[("TECHNICAL PROPOSAL", 14, True, FIRE)]])
textbox(s, Inches(0.7), Inches(1.28), Inches(8.9), Inches(1.4),
        [[("PROTON Customer Complaint", 32, True, WHITE)],
         [("Management System", 32, True, WHITE)]],
        line_spacing=1.0, space_after=2)
textbox(s, Inches(0.7), Inches(2.74), Inches(8.7), Inches(0.5),
        [[("Self-managed, AI-native CRM on Google Cloud", 18, True, POPPY)]])
textbox(s, Inches(0.7), Inches(3.4), Inches(8.6), Inches(0.5),
        [[("Prepared for  ", 15, False, PINK), ("PROTON — e.MAS Customer Operations", 15, True, WHITE)]])
textbox(s, Inches(0.7), Inches(3.9), Inches(8.6), Inches(0.4),
        [[("Prepared by Devoteam  ·  built to your 5-channel complaint SOP", 12, False, WHITE)]])
textbox(s, Inches(0.7), Inches(4.33), Inches(8.6), Inches(0.4),
        [[("July 2026   ·   Version 1.1   ·   Demonstrated live 28 July 2026   ·   Confidential", 10, False, LIGHT)]])
logo(s, dark=False, w=Inches(1.6), x=Inches(0.68), y=SH - Inches(0.92))
npg()

# ---------- 2. EXECUTIVE SUMMARY ----------
s = slide()
page_header(s, "Section 1", "Executive Summary")
textbox(s, Inches(0.55), Inches(1.56), Inches(5.35), Inches(3.5),
        [[("The situation.  ", 13, True, POPPY),
          ("PROTON's complaint operation runs today on Zendesk — a per-seat SaaS CRM. Cost "
           "scales with every agent and premium feature, and the data and AI behaviour live inside "
           "a vendor you don't control.", 13, False, DARK)],
         [("The proposal.  ", 13, True, POPPY),
          ("Replace Zendesk with a self-managed CRM — Chatwoot Community plus a Google Gemini / "
           "Vertex AI layer on Google Cloud — built directly to your 5-channel complaint SOP. The "
           "platform was demonstrated live to your team on 28 July 2026: this is a migration and "
           "hardening engagement, not a greenfield build.", 13, False, DARK)]],
        line_spacing=1.16, space_after=10)
rect(s, Inches(6.2), Inches(1.56), Inches(3.25), Inches(3.6), AQUA)
rect(s, Inches(6.2), Inches(1.56), Inches(3.25), Inches(0.5), DARK)
textbox(s, Inches(6.4), Inches(1.64), Inches(3.0), Inches(0.4),
        [[("WHY NOW / WHY US", 12, True, WHITE)]])
bullets(s, Inches(6.4), Inches(2.24), Inches(2.9), Inches(2.85),
        [("Cost — ", "no per-agent SaaS licensing; pay for infrastructure, not seats"),
         ("Ownership — ", "conversations, KB, AI prompts & data inside PROTON's GCP tenancy"),
         ("Proven — ", "demonstrated live against your SOP on 28 Jul 2026"),
         ("AI-native — ", "Gemini answers in the customer's language, grounded on your FAQ/KB")],
        size=11, gap=9, mcolor=POPPY)
footer(s, npg())

# ---------- 3. YOUR 6 GOALS ----------
s = slide()
page_header(s, "Section 1", "What You Set Out to Achieve")
textbox(s, Inches(0.55), Inches(1.48), Inches(8.9), Inches(0.3),
        [[("Your six goals for the Customer Complaint Management System — taken directly from your "
           "Business Requirement.", 10, False, MIDGREY, True)]])
goals = [
    ("Unified Omni-Channel Management", "Call, Email, WhatsApp & Social in one agent view, single customer record"),
    ("Smart Task Assignment & Reminders", "Status-aware auto-assign by channel priority; sound / desktop / in-app reminders & timeouts"),
    ("Built-in Intelligent FAQ", "Real-time editable KB; keyword-matched suggested replies; one-click reference & quality scoring"),
    ("Automated Workflow & Escalation", "Rule-engine → PIC by category + SOP; email CC + WA alert; 8h auto-escalate / 48h alarm"),
    ("Real-Time Data & Performance", "Dashboards for complaint trends, team performance, SLA & CSAT; anomaly alerts"),
    ("Flexible & Scalable Design", "Grows to new channels, locations & business units; multi-tenant by design"),
]
_chip = [POPPY, MINT, BLUE, FIRE, POPPY, MINT]
_tc = [WHITE, WHITE, WHITE, DARK, WHITE, WHITE]
cardw = Inches(4.35); cardh = Inches(1.0)
col_x = [Inches(0.55), Inches(5.1)]
y0 = int(Inches(1.82)); step = int(Inches(1.11))
for i, (t, d) in enumerate(goals):
    col = i % 2; row = i // 2
    x = col_x[col]; y = Emu(y0 + step*row)
    rect(s, x, y, cardw, cardh, LIGHT)
    rect(s, x, y, Inches(0.09), cardh, _chip[i])
    rect(s, x + Inches(0.22), Emu(int(y) + int(Inches(0.19))), Inches(0.44), Inches(0.44),
         _chip[i], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, x + Inches(0.22), Emu(int(y) + int(Inches(0.19))), Inches(0.44), Inches(0.44),
            [[(str(i+1), 16, True, _tc[i])]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, x + Inches(0.8), Emu(int(y) + int(Inches(0.13))), Inches(3.4), Inches(0.3),
            [[(t, 11.5, True, DARK)]])
    textbox(s, x + Inches(0.8), Emu(int(y) + int(Inches(0.42))), Inches(3.42), Inches(0.5),
            [[(d, 8.5, False, MIDGREY)]], line_spacing=1.0)
footer(s, npg())

# ---------- 4. REQUIREMENT CLUSTERS → PLATFORM ----------
s = slide()
page_header(s, "Section 1", "Requirements → the Platform")
textbox(s, Inches(0.55), Inches(1.5), Inches(8.9), Inches(0.3),
        [[("Every requirement cluster from your BRD and Process-Flow workbook maps to a built "
           "capability.", 10, False, MIDGREY, True)]])
rows = [["#", "Requirement cluster", "What you need"],
        ["1", "Omnichannel inbound", "Call / Email / WhatsApp / Social in one view, single record, new-message alerts, call voice-to-text"],
        ["2", "Agent management", "On-duty check, per-agent channel priorities, status-aware polling assignment, reminders & timeouts"],
        ["3", "FAQ + AI support", "Live-editable KB, keyword-matched suggested replies, one-click reference, FAQ quality scoring"],
        ["4", "Escalation & SLA", "Rule-engine → PIC by category/SOP, email CC + WA alert, WIP/Resolved, 8h/48h auto-escalation"],
        ["5", "Customer 360", "DMS + TSP two-way, auto-identify by number, 4-section 360 card, sync < 3 seconds"],
        ["6", "Reporting / BI", "Channel & division analytics, dept/PIC & Call-Centre KPIs, NPS, CRR, anomaly alerts, Power BI"]]
make_table(s, Inches(0.55), Inches(1.82), Inches(8.9), len(rows),
           [Inches(0.4), Inches(2.1), Inches(6.4)], rows,
           row_h=Inches(0.44), header_h=Inches(0.3), font=9, zebra_fill=AQUA)
textbox(s, Inches(0.55), Inches(4.88), Inches(8.9), Inches(0.3),
        [[("Cross-cutting:  ", 9, True, POPPY),
          ("reduce SaaS cost · own the data & AI · honour the 5-channel SOP timers & surveys exactly · multi-tenant.",
           9, False, DARK)]])
footer(s, npg())

# ---------- 5. 5-CHANNEL SOP AT A GLANCE ----------
s = slide()
page_header(s, "Section 1", "Your 5-Channel SOP at a Glance")
textbox(s, Inches(0.55), Inches(1.48), Inches(8.9), Inches(0.3),
        [[("Business hours: Mon–Fri 08:30–17:30 · Sat/Sun & PH 09:00–17:00. Timers below are taken "
           "verbatim from your Process-Flow workbook.", 9.5, False, MIDGREY, True)]])
rows = [["Channel", "Entry / trigger", "Key SLA & timers", "Resolution & routing"],
        ["WhatsApp", "AI disclaimer;\nsame-language", "Idle 10 min → warn; auto-close 10 / 15 min; agent ack ≤ 2 min", "YES/NO + AI rating; unresolved → agent next business hour"],
        ["Social (FB/IG)", "Post → log ticket ID", "Ack ≤ 2 working hours", "Social = 1st priority when opted; next business hour"],
        ["Email", "e.mascentre@pronet.my", "Auto-ack once per new thread; status ≤ 4 working hours", "1 business-day response; Escalation Policy"],
        ["IVR Call", "1300-888-877;\nfemale AI voice", "Office-hour check; agent ≤ 20 sec; queue-busy prompt (EN + BM)", "RSA 24h vs office-hours; escalate to team"],
        ["SSI Survey", "e.MAS app (day 11,\nexpires day 14)", "Target > 90%; response rate > 45%", "Ingest & report only — lives in the e.MAS app"]]
make_table(s, Inches(0.55), Inches(1.82), Inches(8.9), len(rows),
           [Inches(1.35), Inches(1.9), Inches(3.15), Inches(2.5)], rows,
           row_h=Inches(0.62), header_h=Inches(0.3), font=8.5, zebra_fill=BLUELT)
footer(s, npg())

# ---------- 6. WHAT TO EXPECT ----------
s = slide()
page_header(s, "Section 1", "What to Expect")
rows = [["Business Expectation", "Technical Expectation"],
        ["No per-agent / per-seat licence fees", "Chatwoot Community (open-source), self-hosted on GCP — unlimited agents"],
        ["Your data stays yours", "All conversations, KB and AI prompts inside PROTON's GCP project"],
        ["AI answers instantly, in their language", "Gemini on Vertex AI, same-language replies, grounded on your FAQ/KB"],
        ["Your SOP, not a generic product", "Process flow as code: disclaimers, idle-close, YES/NO resolution, rating surveys"],
        ["Faster complaint resolution", "Auto-classification → PIC routing, SLA timers (2-min WA ack … 48h alarm)"],
        ["Management visibility", "BigQuery analytics + Looker / Power BI, scheduled PDF/Excel, anomaly alerts"],
        ["Predictable, scalable cost", "Managed GCP services; pay for compute/storage consumed, scale per tenant"],
        ["Low delivery risk & continuity", "Majority already built & demoed live; phased cutover, parallel-run first"]]
make_table(s, Inches(0.55), Inches(1.55), Inches(8.9), len(rows),
           [Inches(3.9), Inches(5.0)], rows,
           row_h=Inches(0.4), header_h=Inches(0.34), font=10.5)
footer(s, npg())

# ---------- 7. DIVIDER §2 ----------
divider(2, "What We've Built & Demonstrated", "Shown live to your team · 28 July 2026")

# ---------- 8. LIVE DEMO RECAP ----------
s = slide()
page_header(s, "Section 2", "Demonstrated Live — Working Today")
textbox(s, Inches(0.55), Inches(1.62), Inches(4.4), Inches(0.3), [[("AGENT WORKSPACE", 12, True, POPPY)]])
bullets(s, Inches(0.55), Inches(2.0), Inches(4.4), Inches(2.9),
        [("Unified inbox ", "— one agent view across channels"),
         ("AI suggest-reply ", "— drafted from your knowledge base"),
         ("Summarize ", "— conversation → private note for the agent"),
         ("Ask Copilot ", "— grounded on internal / customer data"),
         ("FAQ assist ", "— keyword-matched replies, one click")],
        size=11.5, gap=10)
textbox(s, Inches(5.15), Inches(1.62), Inches(4.3), Inches(0.3), [[("RECORDS & REPORTING", 12, True, POPPY)]])
bullets(s, Inches(5.15), Inches(2.0), Inches(4.3), Inches(2.9),
        [("Contact panel ", "— info, history, edit & merge contacts"),
         ("Conversation history ", "— prior chats per customer"),
         ("Native reports ", "— overview, open conversations, agent status"),
         ("SLA metrics ", "— avg resolution ~6m30s, per-channel volume"),
         ("Auto-labels ", "— AI classifies inquiry / complaint / feedback")],
        size=11.5, gap=10, mcolor=MINT)
rect(s, Inches(0.55), Inches(4.72), Inches(8.9), Inches(0.5), AQUA)
textbox(s, Inches(0.72), Inches(4.72), Inches(8.6), Inches(0.5),
        [[("All shown live on your PROTON tenant — across WhatsApp, Website, API and Voice.",
           10.5, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, npg())

# ---------- 9. AI & KNOWLEDGE ----------
s = slide()
page_header(s, "Section 2", "AI & Knowledge — Under Your Control")
textbox(s, Inches(0.55), Inches(1.62), Inches(4.4), Inches(0.3), [[("AI IN THE CONVERSATION", 12, True, POPPY)]])
bullets(s, Inches(0.55), Inches(2.0), Inches(4.4), Inches(2.9),
        [("Suggest reply & FAQ assist ", "— KB-grounded, same knowledge base"),
         ("Summarize ", "— long threads into an agent private note"),
         ("Ask Copilot ", "— grounded on your internal data"),
         ("Auto-classify ", "— inquiry / complaint / feedback labels"),
         ("Same-language ", "— answers in the customer's language")],
        size=11.5, gap=10)
textbox(s, Inches(5.15), Inches(1.62), Inches(4.3), Inches(0.3), [[("KNOWLEDGE & GUARDRAILS", 12, True, POPPY)]])
bullets(s, Inches(5.15), Inches(2.0), Inches(4.3), Inches(2.9),
        [("Knowledge base ", "— manual FAQ + website ingest → Vertex AI Search"),
         ("Playground ", "— test answers before they go live"),
         ("Custom tools ", "— call your DMS / internal APIs"),
         ("Persona ", "— per-assistant instruction, temperature, guardrails"),
         ("Model ", "— Gemini 2.5 Flash, configurable")],
        size=11.5, gap=10, mcolor=MINT)
rect(s, Inches(0.55), Inches(4.72), Inches(8.9), Inches(0.5), BEIGE)
textbox(s, Inches(0.72), Inches(4.72), Inches(8.6), Inches(0.5),
        [[("Operators own the FAQ, the AI prompts and the guardrails — no black box, no vendor lock-in.",
           10.5, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, npg())

# ---------- 10. OMNICHANNEL — LIVE CHANNEL STATUS ----------
s = slide()
page_header(s, "Section 2", "Omnichannel — Live Channel Status")
rows = [["Channel", "Integration", "Status (as demonstrated 28 Jul)"],
        ["WhatsApp", "Twilio", "● Live — demonstrated"],
        ["Website widget", "Native", "● Live — demonstrated"],
        ["API (backend-to-backend)", "PROTON API", "● Live — demonstrated"],
        ["Voice", "Twilio conversational AI", "● Live — real call demonstrated"],
        ["Email", "IMAP / SMTP", "◐ Pending — PROTON domain + credentials"],
        ["Social (FB / IG)", "Meta Business", "◐ Pending — Meta business verification"],
        ["Telegram / Line", "Native connectors", "○ Available on request"]]
make_table(s, Inches(0.55), Inches(1.6), Inches(8.9), len(rows),
           [Inches(2.7), Inches(2.9), Inches(3.3)], rows,
           row_h=Inches(0.4), header_h=Inches(0.34), font=10.5, zebra_fill=MINTLT)
textbox(s, Inches(0.55), Inches(5.04), Inches(8.9), Inches(0.28),
        [[("Multimodal on WhatsApp (voice note / image / video) unlocks once Meta business verification is complete.",
           9, False, MIDGREY, True)]])
footer(s, npg())

# ---------- 11. CONVERSATIONAL VOICE AI ----------
s = slide()
page_header(s, "Section 2", "Conversational Voice AI (Twilio)")
flow = [("Call 1300-888-877", "female AI voice", MIDGREY),
        ("LLM understands", "no press-1 / press-2 menu", FIRE),
        ("KB-grounded answer", "X70 specs · AEB/FCW/ACC/LKA", BLUE),
        ("CSAT rating", "rate 1–5", MINT)]
fw = Inches(2.0); fh = Inches(0.82); fy = Inches(1.78); gap = Inches(0.24); fx = Inches(0.55)
for i, (t, d, c) in enumerate(flow):
    chip(s, fx, fy, fw, fh, t, d, c, tcolor=(DARK if c is FIRE else WHITE),
         scolor=(DARK if c is FIRE else WHITE), tsize=10.5, ssize=8)
    if i < len(flow) - 1:
        ax0 = Emu(int(fx) + int(fw))
        connector(s, ax0, fy + fh/2, Emu(int(ax0) + int(gap)), fy + fh/2, color=POPPY)
    fx = Emu(int(fx) + int(fw) + int(gap))
# handoff row
hy = Inches(3.1)
chip(s, Inches(0.55), hy, Inches(3.0), Inches(0.7), "Escalate to human",
     "orchestrator — agent-B layer", DARK, tsize=11, ssize=8)
connector(s, Inches(3.55), hy + Inches(0.35), Inches(4.15), hy + Inches(0.35), color=POPPY)
routes = [("Sales", Inches(4.25)), ("Repair", Inches(6.05)), ("RSA — 24h", Inches(7.85))]
for name, x in routes:
    chip(s, x, hy, Inches(1.55), Inches(0.7), name, "", BLUE, tsize=11)
textbox(s, Inches(0.55), Inches(4.12), Inches(8.9), Inches(0.9),
        [[("How it works.  ", 11, True, POPPY),
          ("The LLM understands intent directly, so callers skip the IVR menu. AI answers vehicle "
           "questions 24/7; live-agent handoff routes by context to Sales, Repair or 24h Roadside "
           "Assistance, with business-hour rules.", 11, False, DARK)]], line_spacing=1.15)
textbox(s, Inches(0.55), Inches(5.02), Inches(8.9), Inches(0.3),
        [[("Decision to confirm:  ", 9, True, POPPY),
          ("conversational LLM-native voice vs. Twilio IVR (press-1 / press-2) for production.", 9, False, MIDGREY, True)]])
footer(s, npg())

# ---------- 12. LIFECYCLE & SLA MAPPED TO SOP ----------
s = slide()
page_header(s, "Section 2", "Lifecycle & SLA — Built to Your SOP")
steps = [("1. Disclaimer & same-language", "AI greets, answers in the customer's language, grounded on your FAQ."),
         ("2. Idle handling", "10-min inactivity warning (\"chat closes in 5 min\"); auto-close after 10 / 15 min."),
         ("3. Resolution gate", "\"Is your case resolved? YES / NO\" + AI-performance rating survey."),
         ("4. Escalate to agent", "Unresolved → active agent by priority (WA → Call → Email → Social); ack ≤ 2 min."),
         ("5. Categorize", "AI assigns case category / division from your predefined list on close."),
         ("6. PIC escalation", "Rule-engine → PIC by category + SOP; CC email (+ optional WhatsApp alert)."),
         ("7. Time-based alarms", "8h auto-escalate to higher level; 48h unresolved alarm; WIP/Resolved recorded (agent · time · remarks).")]
y = Inches(1.6)
for lead, desc in steps:
    rect(s, Inches(0.55), y+Inches(0.02), Inches(0.14), Inches(0.42), FIRE)
    textbox(s, Inches(0.82), y, Inches(8.6), Inches(0.5),
            [[(lead + "  —  ", 11.5, True, POPPY), (desc, 11.5, False, DARK)]], line_spacing=1.0)
    y += Inches(0.5)
footer(s, npg())

# ---------- 13. DIVIDER §3 ----------
divider(3, "Target Architecture", "Zendesk (SaaS)  →  self-managed on Google Cloud")

# ---------- 14. CURRENT vs TARGET ----------
s = slide()
page_header(s, "Section 3", "Current State vs. Target State")
bx = Inches(0.55); bw = Inches(3.55); by = Inches(1.72); bh = Inches(3.05)
rect(s, bx, by, bw, bh, LIGHT)
rect(s, bx, by, bw, Inches(0.44), MIDGREY)
textbox(s, bx, by + Inches(0.06), bw, Inches(0.34), [[("TODAY — ZENDESK (SaaS)", 12, True, WHITE)]], align=PP_ALIGN.CENTER)
bullets(s, bx + Inches(0.22), by + Inches(0.6), bw - Inches(0.44), Inches(2.4),
        ["Per-agent subscription; SLA / roles / AI behind paid tiers",
         "Data & AI behaviour hosted by the vendor",
         "FAQ / Guide inside Zendesk — limited AI-grounding control",
         "DMS/TSP & telephony bolted onto a closed platform"],
        size=10.5, gap=9, marker="✕", mcolor=MIDGREY)
midx = Inches(4.35)
rect(s, midx, Inches(2.86), Inches(0.7), Inches(0.78), POPPY, shape=MSO_SHAPE.RIGHT_ARROW)
textbox(s, Inches(4.2), Inches(2.5), Inches(1.0), Inches(0.3), [[("MIGRATE", 9.5, True, POPPY)]], align=PP_ALIGN.CENTER)
ax = Inches(5.35); aw = Inches(4.1)
rect(s, ax, by, aw, bh, AQUA)
rect(s, ax, by, aw, Inches(0.44), DARK)
textbox(s, ax, by + Inches(0.06), aw, Inches(0.34), [[("TARGET — SELF-MANAGED ON GCP", 12, True, WHITE)]], align=PP_ALIGN.CENTER)
bullets(s, ax + Inches(0.22), by + Inches(0.6), aw - Inches(0.44), Inches(2.4),
        ["Open-source Chatwoot core — no seat licences; SLA / roles / AI built, not rented",
         "Runs inside PROTON's GCP project — full data residency & control",
         "KB + AI prompts owned and edited by PROTON operators (no-code)",
         "Native first-party integration surfaces for DMS/TSP, telephony & BI"],
        size=10.5, gap=9, marker="✓", mcolor=MINT)
textbox(s, Inches(0.55), Inches(4.95), Inches(8.9), Inches(0.35),
        [[("Same capabilities — moved onto foundations PROTON owns and controls, at infrastructure cost.",
           10.5, True, DARK)]], align=PP_ALIGN.CENTER)
footer(s, npg())

# ---------- 15. GCP PRODUCTION ARCHITECTURE ----------
s = slide()
page_header(s, "Section 3", "GCP Production Architecture")
textbox(s, Inches(0.55), Inches(1.62), Inches(1.0), Inches(0.4), [[("Internet", 10, True, MIDGREY)]], anchor=MSO_ANCHOR.MIDDLE)
lb_x, lb_y, lb_w, lb_h = Inches(1.55), Inches(1.6), Inches(1.65), Inches(0.5)
chip(s, lb_x, lb_y, lb_w, lb_h, "Cloud LB + Caddy", "TLS · routing · per-tenant entry", POPPY, tsize=10.5, ssize=8)
connector(s, Inches(1.42), Inches(1.85), lb_x, Inches(1.85), color=MIDGREY)
app_y = Inches(2.45); app_h = Inches(0.62); app_w = Inches(2.7)
apps = [("Chatwoot  (GKE)", "Rails + Sidekiq · CRM / live-chat core", Inches(0.75)),
        ("agent  (GKE)", "Webhook sync · AI orchestration · SLA", Inches(3.6)),
        ("backend  (GKE)", "Gemini agent · KB · routing · metrics", Inches(6.45))]
rect(s, Inches(0.55), Inches(2.3), Inches(8.9), Inches(0.94), MINTLT)
textbox(s, Inches(0.62), Inches(2.32), Inches(5.2), Inches(0.2), [[("APPLICATION — GKE Standard  (containerized · node-pool control · multi-tenant)", 8, True, MINT)]])
for t, d, x in apps:
    chip(s, x, app_y, app_w, app_h, t, d, MINT, tsize=10.5, ssize=7.5)
    connector(s, x + app_w/2, lb_y + lb_h, x + app_w/2, app_y, color=MIDGREY, arrow=False, width=1.25)
data_y = Inches(3.6); data_h = Inches(0.66); dw = Inches(2.12)
rect(s, Inches(0.55), Inches(3.46), Inches(8.9), Inches(0.98), BLUELT)
textbox(s, Inches(0.62), Inches(3.48), Inches(4.0), Inches(0.2), [[("AI · DATA · STATE  (managed GCP services)", 8, True, BLUE)]])
tiles = [("Vertex AI — Gemini", "drafts · classify · same-language", BLUE),
         ("Vertex Search + pgvector", "KB grounding — no black box", BLUE),
         ("Cloud SQL (Postgres)", "Chatwoot DB + per-tenant KB", DARK),
         ("Memorystore · Storage", "Redis queues · attachments", DARK)]
tx = Inches(0.62)
for t, d, c in tiles:
    chip(s, tx, data_y, dw, data_h, t, d, c, tsize=9.5, ssize=7.5)
    tx = Emu(int(tx) + int(dw) + int(Inches(0.13)))
rect(s, Inches(0.55), Inches(4.62), Inches(8.9), Inches(0.5), BEIGE)
textbox(s, Inches(0.72), Inches(4.62), Inches(8.6), Inches(0.5),
        [[("Secret Manager  ·  BigQuery → Looker Studio / Power BI  ·  Cloud Monitoring & Logging  ·  "
           "Automated backups  ·  Multi-tenant isolation (per-tenant namespaces + node pools · isolated DBs)",
           9, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, npg())

# ---------- 16. DATA MIGRATION ----------
s = slide()
page_header(s, "Section 3", "Data Migration — Zendesk → Platform")
flow = [("Zendesk API", "export", MIDGREY), ("Transform &", "map", FIRE),
        ("Bulk import", "Chatwoot API", MINT), ("KB ingest", "Vertex + pgvector", BLUE)]
fw = Inches(1.9); fh = Inches(0.8); fy = Inches(1.85); gap = Inches(0.32); fx = Inches(0.55)
for i, (t, d, c) in enumerate(flow):
    chip(s, fx, fy, fw, fh, t, d, c, tcolor=(DARK if c is FIRE else WHITE),
         scolor=(DARK if c is FIRE else WHITE), tsize=11, ssize=9)
    if i < len(flow) - 1:
        ax0 = Emu(int(fx) + int(fw))
        connector(s, ax0, fy + fh/2, Emu(int(ax0) + int(gap)), fy + fh/2, color=POPPY)
    fx = Emu(int(fx) + int(fw) + int(gap))
textbox(s, Inches(0.55), Inches(3.05), Inches(4.4), Inches(0.3), [[("WHAT WE MIGRATE", 12, True, POPPY)]])
bullets(s, Inches(0.55), Inches(3.45), Inches(4.4), Inches(1.7),
        [("Tickets / conversations ", "→ Chatwoot conversations (status & labels)"),
         ("Contacts / customers ", "→ Chatwoot contacts (phone/email, custom fields)"),
         ("Knowledge / Guide (FAQ) ", "→ PROTON-owned KB (Vertex Search + pgvector)"),
         ("Attachments → Storage ", "· reporting history → BigQuery")],
        size=10.5, gap=7)
textbox(s, Inches(5.15), Inches(3.05), Inches(4.3), Inches(0.3), [[("HOW — DE-RISKED, NO BIG-BANG", 12, True, POPPY)]])
bullets(s, Inches(5.15), Inches(3.45), Inches(4.3), Inches(1.7),
        [("Export → transform ", "→ bulk import via Chatwoot API + KB ingest"),
         ("Parallel run ", "— platform live alongside Zendesk during validation"),
         ("Cutover ", "— switch channel endpoints once data fidelity validated"),
         ("Decommission ", "Zendesk after sign-off")],
        size=10.5, gap=7, mcolor=MINT)
footer(s, npg())

# ---------- 17. DIVIDER §4 ----------
divider(4, "Your Feedback → What We Build Next", "Priorities captured in the 28 July session")

# ---------- 18. GAP-CLOSURE PRIORITIES ----------
s = slide()
page_header(s, "Section 4", "Gap-Closure Priorities")
rows = [["Pri.", "Item", "What you asked for", "Dependency"],
        ["P1", "Email escalation", "Auto-ack once per thread + forward internal case as a separate conversation (no CC/BCC), SOP-driven", "PROTON email domain + IMAP/SMTP"],
        ["P1", "Customer 360", "Single customer ID (name / vehicle / phone); group all conversations; 4-section card < 3s", "DMS + TSP API access"],
        ["P1", "Case management", "Per-ticket support dashboard + agent \"my cases\" view; hierarchical categories (main + sub)", "—"],
        ["P2", "Reporting & BI", "Customization + Power BI (BigQuery embed); role-based visibility (admin/leader vs agent)", "PROTON report examples"],
        ["P2", "Agent status", "Auto-busy on call; round-robin skips busy; configurable concurrent-ticket limits", "Telephony events"],
        ["P2", "Multimodal + language", "WhatsApp voice-note / image / video; always answer in customer's language; bulk FAQ/PDF upload", "Meta verification"]]
make_table(s, Inches(0.55), Inches(1.55), Inches(8.9), len(rows),
           [Inches(0.55), Inches(1.75), Inches(4.55), Inches(2.05)], rows,
           row_h=Inches(0.56), header_h=Inches(0.3), font=8.8, zebra_fill=POPPYLT)
footer(s, npg())

# ---------- 19. EMAIL ESCALATION FLOW ----------
s = slide()
page_header(s, "Section 4", "Email Escalation — Your Priority Flow")
# row 1: customer email -> CRM inbox
chip(s, Inches(0.55), Inches(1.7), Inches(2.3), Inches(0.72), "Customer email", "to e.mascentre@pronet.my", MIDGREY, tsize=11, ssize=8)
connector(s, Inches(2.85), Inches(2.06), Inches(3.25), Inches(2.06), color=POPPY)
chip(s, Inches(3.25), Inches(1.7), Inches(2.3), Inches(0.72), "CRM inbox", "IMAP / SMTP", BLUE, tsize=11, ssize=8)
# branch up: auto-ack (customer)
connector(s, Inches(5.55), Inches(2.06), Inches(6.15), Inches(2.06), color=POPPY)
chip(s, Inches(6.15), Inches(1.62), Inches(3.3), Inches(0.9), "Auto-ack to customer — ONCE",
     "new thread only · not on replies / agent sends", MINT, tsize=10, ssize=8)
# branch down from CRM inbox → routed into "Agent creates case" (Z-path)
connector(s, Inches(4.4), Inches(2.42), Inches(4.4), Inches(2.68), color=POPPY, arrow=False)
connector(s, Inches(4.4), Inches(2.68), Inches(1.825), Inches(2.68), color=POPPY, arrow=False)
connector(s, Inches(1.825), Inches(2.68), Inches(1.825), Inches(2.95), color=POPPY)
chip(s, Inches(0.55), Inches(2.95), Inches(2.55), Inches(0.78), "Agent creates case", "classify by category", DARK, tsize=10.5, ssize=8)
connector(s, Inches(3.1), Inches(3.34), Inches(3.5), Inches(3.34), color=POPPY)
chip(s, Inches(3.5), Inches(2.95), Inches(2.55), Inches(0.78), "PIC by category + SOP", "Escalation Policy", FIRE, tcolor=DARK, scolor=DARK, tsize=10.5, ssize=8)
connector(s, Inches(6.05), Inches(3.34), Inches(6.45), Inches(3.34), color=POPPY)
chip(s, Inches(6.45), Inches(2.95), Inches(3.0), Inches(0.78), "Forward internal / dealer", "SEPARATE conversation", BLUE, tsize=10.5, ssize=8)
# rules strip
rect(s, Inches(0.55), Inches(4.05), Inches(8.9), Inches(1.02), AQUA)
textbox(s, Inches(0.75), Inches(4.14), Inches(8.5), Inches(0.3), [[("BUILT TO YOUR SOP", 11, True, POPPY)]])
bullets(s, Inches(0.75), Inches(4.46), Inches(8.5), Inches(0.6),
        [("Two separate emails — ", "customer acknowledgement and internal case; no CC / BCC / thread trailing."),
         ("Status ≤ 4 working hours ", "· WIP/Resolved · 8h auto-escalate / 48h alarm · PROTON supplies email service + subdomain + credentials.")],
        size=9.5, gap=5, mcolor=DARK)
footer(s, npg())

# ---------- 20. CUSTOMER 360 & CASE MANAGEMENT ----------
s = slide()
page_header(s, "Section 4", "Customer 360 & Case Management")
rect(s, Inches(0.55), Inches(1.62), Inches(4.35), Inches(3.5), BLUELT)
textbox(s, Inches(0.75), Inches(1.72), Inches(4.0), Inches(0.3), [[("CUSTOMER 360 VIEW CARD", 11.5, True, BLUE)]])
bullets(s, Inches(0.75), Inches(2.12), Inches(3.95), Inches(2.9),
        [("Auto-pop ", "on call / WA by caller or WA number"),
         ("Personal ", "— name, gender, contact, address"),
         ("Vehicle ", "— VRN, VIN, model, dealer, purchase date"),
         ("Service history ", "— recent service, dealer, open RO status"),
         ("Call-centre history ", "— Call / WA / email / social / RSA"),
         ("DMS + TSP ", "two-way · sync < 3s, async loading")],
        size=10.5, gap=7, mcolor=BLUE)
rect(s, Inches(5.1), Inches(1.62), Inches(4.35), Inches(3.5), POPPYLT)
textbox(s, Inches(5.3), Inches(1.72), Inches(4.0), Inches(0.3), [[("CASE MANAGEMENT", 11.5, True, POPPY)]])
bullets(s, Inches(5.3), Inches(2.12), Inches(3.95), Inches(2.9),
        [("Single customer ID ", "(CIF-style) groups all conversations by name / vehicle / phone"),
         ("Hierarchical categories ", "— main + dependent subcategory"),
         ("Priority ", "— high / medium / low, assign agent or team"),
         ("Per-ticket dashboard ", "— caller, number, status at a glance"),
         ("My cases ", "— agent follow-up view of open tickets"),
         ("States ", "— WIP / Resolved / Temporary-Closed / Closed")],
        size=10.5, gap=7, mcolor=POPPY)
footer(s, npg())

# ---------- 21. REPORTING & BI ----------
s = slide()
page_header(s, "Section 4", "Reporting & Business Intelligence")
textbox(s, Inches(0.55), Inches(1.6), Inches(4.5), Inches(0.3), [[("REPORTS YOU ASKED FOR", 12, True, POPPY)]])
bullets(s, Inches(0.55), Inches(1.98), Inches(4.5), Inches(3.1),
        [("Channel source ", "— Call / WA / Email / Social"),
         ("Case division ", "— Apps / Sales / Aftersales / Charging"),
         ("Dept & PIC ", "— first-response, resolution, CRR ranking"),
         ("Call-Centre KPI ", "— SLA achievement, tasks/agent, closure time"),
         ("NPS ", "— customer rates agent on Call & WA"),
         ("Case lifecycle ", "— Higher-escalation / WIP / Temp-Closed / Closed"),
         ("Trends ", "— daily / weekly / monthly, peak-hour analysis")],
        size=10.5, gap=7)
rect(s, Inches(5.15), Inches(1.62), Inches(4.3), Inches(3.05), MINTLT)
textbox(s, Inches(5.35), Inches(1.74), Inches(4.0), Inches(0.3), [[("HOW WE DELIVER IT", 12, True, MINT)]])
bullets(s, Inches(5.35), Inches(2.14), Inches(3.9), Inches(2.5),
        [("Native reports ", "in the CRM today (overview, agent status, SLA)"),
         ("BigQuery warehouse ", "→ Power BI / Looker (connect or embed)"),
         ("Scheduled exports ", "— PDF / Excel auto-sent to management"),
         ("Anomaly dashboard ", "— real-time channel-spike alerts"),
         ("Role-based views ", "— admin / leader vs agent")],
        size=10.5, gap=8, mcolor=MINT)
textbox(s, Inches(0.55), Inches(5.02), Inches(8.9), Inches(0.3),
        [[("Next:  ", 9, True, POPPY), ("PROTON to share report / visualization examples so we can tailor the dashboards.", 9, False, MIDGREY, True)]])
footer(s, npg())

# ---------- 22. DIVIDER §5 ----------
divider(5, "Delivery & Commercials", "Timeline · Scope · Deliverables · Engagement")

# ---------- 23. PROJECT TIMELINE ----------
s = slide()
page_header(s, "Section 5", "Project Timeline — ~12–16 Weeks, Phased")
rows = [["Phase", "Weeks", "Focus", "Key outcomes"],
        ["P0 — Discovery &\nGCP foundation", "1–2", "Access, data contracts, GCP project / landing zone",
         "GCP org/project, network, Secret Manager, CI; Zendesk export sample validated"],
        ["P1 — Core platform\ncutover", "3–7", "Stand up production stack; migrate data; wire channels",
         "Chatwoot + agent + backend on GKE + Cloud SQL; WA/Web/API live; KB migrated; AI + SOP flows live"],
        ["P2 — Gap closure", "6–12", "The 28-Jul priorities (parallel with P1 tail)",
         "Email escalation; Customer 360 + DMS/TSP; per-ticket dashboard + hierarchical categories; Reporting/Power BI"],
        ["P3 — Hardening,\nparallel-run & handover", "12–16", "HA, monitoring, UAT, cutover, decommission",
         "Autoscaling/HA, backups, dashboards; UAT sign-off; Zendesk decommissioned; runbook + training"]]
make_table(s, Inches(0.55), Inches(1.6), Inches(8.9), len(rows),
           [Inches(1.75), Inches(0.7), Inches(2.85), Inches(3.6)], rows,
           row_h=Inches(0.74), header_h=Inches(0.32), font=9.5, zebra_fill=BLUELT)
textbox(s, Inches(0.55), Inches(5.02), Inches(8.9), Inches(0.3),
        [[("P0/P1 largely assemble already-built components; the real engineering is P2's external-dependency items (email, DMS/TSP, telephony).",
           9, False, MIDGREY, True)]])
footer(s, npg())

# ---------- 24. SCOPE OF WORK ----------
s = slide()
page_header(s, "Section 5", "Scope of Work")
sc1 = ["GCP landing zone: project, networking, IAM, Secret Manager, CI/CD",
       "Deploy production Chatwoot + agent + backend on GKE Standard + Cloud SQL",
       "Omnichannel: WhatsApp, Email, Social (FB/IG), Web, Voice (conversational)",
       "AI layer: Gemini reply drafting, same-language, KB grounding, auto-classification",
       "Knowledge base: migration + no-code authoring; bulk FAQ / PDF ingest",
       "Agent management: presence, channel-priority routing, auto-busy, My-Tasks timers"]
sc2 = ["Escalation & SLA: category→PIC, email CC + WA alert, 8h/48h, WIP/Resolved",
       "Email escalation: auto-ack + separate internal case per Escalation Policy SOP",
       "Lifecycle/SOP: disclaimers, idle-close, YES/NO gate, rating surveys, auto-ack",
       "Customer 360: DMS + TSP two-way + 4-section 360 card (subject to API access)",
       "Case mgmt: hierarchical categories, per-ticket dashboard, RBAC roles",
       "Reporting/BI: BigQuery + Power BI, scheduled exports, anomaly alerts, NPS/CRR",
       "Data migration from Zendesk + Ops: monitoring, backups, multi-tenant, training"]
bullets(s, Inches(0.55), Inches(1.55), Inches(4.45), Inches(3.6), sc1, size=10, gap=8)
bullets(s, Inches(5.15), Inches(1.55), Inches(4.3), Inches(3.6), sc2, size=10, gap=8)
textbox(s, Inches(0.55), Inches(5.02), Inches(8.9), Inches(0.3),
        [[("Assumptions:  ", 9, True, POPPY),
          ("PROTON provides GCP billing/org access, DMS/TSP API + data contracts, email domain, Zendesk export, and a telephony/CTI decision.",
           9, False, DARK)]])
footer(s, npg())

# ---------- 25. OUT OF SCOPE ----------
s = slide()
page_header(s, "Section 5", "Out of Scope")
oo1 = ["Telephony/PSTN carrier procurement — we integrate the CTI/provider; the phone-line contract & number provisioning is PROTON's",
       "DMS / TSP source-system changes — we consume their APIs; internal changes are the owners' responsibility",
       "SSI dealer-satisfaction survey workflow — lives in the e.MAS app; the platform ingests & reports SSI, but does not host that flow"]
oo2 = ["Power BI licences — we connect to BigQuery; Microsoft licensing is PROTON's",
       "Meta business verification — required to unlock Social & WhatsApp multimodal; owned by PROTON",
       "Ongoing FAQ/KB content authoring after handover — enabled via no-code tools, owned by PROTON",
       "Non-CRM business applications, custom hardware, and third-party SaaS beyond the GCP stack"]
bullets(s, Inches(0.55), Inches(1.7), Inches(4.45), Inches(3.4), oo1, size=11, gap=15, marker="✕", mcolor=MIDGREY)
bullets(s, Inches(5.15), Inches(1.7), Inches(4.3), Inches(3.4), oo2, size=11, gap=15, marker="✕", mcolor=MIDGREY)
footer(s, npg())

# ---------- 26. DELIVERABLES ----------
s = slide()
page_header(s, "Section 5", "Deliverables")
rows = [["Deliverable", "Description"],
        ["Production GCP environment", "Fully provisioned, IaC-described GCP project running the platform"],
        ["Migrated CRM", "Chatwoot with imported Zendesk tickets, contacts, KB, attachments"],
        ["AI automation layer", "Gemini reply, classification, KB-grounded answers, SOP flows — live"],
        ["Email escalation", "SOP-driven auto-ack + separate internal case forwarding"],
        ["Customer 360 integration", "DMS/TSP two-way sync + 4-section 360 card (pending API access)"],
        ["Case management & RBAC", "Hierarchical categories, per-ticket dashboard, roles & permissions"],
        ["Reporting suite", "BigQuery + Power BI dashboards, scheduled exports, anomaly alerts, NPS/CRR"],
        ["Data-migration report", "Record counts, fidelity validation, parallel-run results"],
        ["Operations runbook & training", "Deploy, backup/restore, monitoring, tenant provisioning; admin & agent training"],
        ["UAT sign-off & cutover", "Validated go-live and Zendesk decommission"]]
make_table(s, Inches(0.55), Inches(1.55), Inches(8.9), len(rows),
           [Inches(3.1), Inches(5.8)], rows,
           row_h=Inches(0.335), header_h=Inches(0.32), font=10.5, zebra_fill=MINTLT)
footer(s, npg())

# ---------- 27. EFFORT & MANDAYS ----------
s = slide()
page_header(s, "Section 5", "Effort Estimate — Mandays")
textbox(s, Inches(0.55), Inches(1.5), Inches(8.9), Inches(0.26),
        [[("Indicative effort for the ~16-week migration + hardening + gap-closure engagement — detail in the commercial proposal.",
           9.5, False, MIDGREY, True)]])
rows = [["Workstream", "Key roles", "MD"],
        ["Project management & governance", "PM", "24"],
        ["Discovery, design & GCP landing zone (GKE Standard)", "SA · DevOps", "18"],
        ["Platform deploy, CI/CD & data migration", "DevOps · BE", "26"],
        ["Omnichannel integration (WA / Web / Email / Social / Voice)", "BE", "24"],
        ["AI layer, KB & lifecycle / SOP automation", "BE", "28"],
        ["Email escalation + Customer 360 + DMS / TSP (P1)", "BE", "34"],
        ["Case management, RBAC & agent management", "BE · FE", "28"],
        ["Reporting & BI (BigQuery + Power BI)", "BI", "16"],
        ["Frontend / Chatwoot fork customization", "FE", "18"],
        ["QA, UAT, hardening, HA, runbook & training", "QA · DevOps", "36"],
        ["Total — indicative delivery effort", "", "252"]]
tbl = make_table(s, Inches(0.55), Inches(1.72), Inches(8.9), len(rows),
                 [Inches(5.6), Inches(2.1), Inches(1.2)], rows,
                 row_h=Inches(0.26), header_h=Inches(0.3), font=9.5, zebra_fill=LIGHT)
for ri in range(len(rows)):
    for p in tbl.cell(ri, 2).text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
_last = len(rows) - 1
for ci in range(3):
    cell = tbl.cell(_last, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = POPPYLT
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True; r.font.color.rgb = DARK
textbox(s, Inches(0.55), Inches(5.0), Inches(8.9), Inches(0.26),
        [[("Team:  ", 8.5, True, POPPY),
          ("PM · Solution Architect / Tech Lead · Backend/AI · Frontend · Data/BI · QA (blended) — cost = mandays × blended day-rate, in the commercial proposal.",
           8.5, False, DARK)]])
footer(s, npg())

# ---------- 28. GCP PRICING & SPECIFICATION ----------
s = slide()
page_header(s, "Section 5", "GCP Pricing & Specification")
textbox(s, Inches(0.55), Inches(1.48), Inches(8.9), Inches(0.26),
        [[("One production tenant · region asia-southeast1 (Singapore) — nearest full-service region; "
           "the Malaysia (KL) region is not yet GA.", 9, False, MIDGREY, True)]])
rows = [["Service", "Specification", "USD / mo"],
        ["GKE Standard — cluster fee", "Regional cluster ($0.10/hr)", "73"],
        ["GKE Standard — nodes", "3 × e2-standard-4 (4 vCPU / 16 GB)", "362"],
        ["Cloud SQL — PostgreSQL (HA)", "2 vCPU / 8 GB · 100 GB SSD · regional HA", "300"],
        ["Memorystore — Redis (HA)", "Standard tier · 4 GB", "181"],
        ["Vertex AI (Gemini 2.5 Flash + Search)", "usage · ~5k conversations / queries mo", "51"],
        ["Cloud Storage", "100 GB standard + egress", "8"],
        ["Networking (Load Balancer + egress)", "1 forwarding rule + ~100 GB", "32"],
        ["Cloud Logging + Monitoring", "~100 GB ingest / mo", "25"],
        ["Secret Mgr · Artifact Reg · Build · BigQuery", "combined", "17"],
        ["Total — indicative, on-demand", "one production tenant", "1,049"]]
tbl = make_table(s, Inches(0.55), Inches(1.76), Inches(8.9), len(rows),
                 [Inches(3.5), Inches(3.9), Inches(1.5)], rows,
                 row_h=Inches(0.255), header_h=Inches(0.3), font=9.5, zebra_fill=BLUELT)
for ri in range(len(rows)):
    for p in tbl.cell(ri, 2).text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
_last = len(rows) - 1
for ci in range(3):
    cell = tbl.cell(_last, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = POPPYLT
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True; r.font.color.rgb = DARK
textbox(s, Inches(0.55), Inches(4.74), Inches(8.9), Inches(0.5),
        [[("1-year committed-use discount ", 9, True, POPPY),
          ("lowers compute to ~$233/mo → total ≈ $920/mo · multi-tenant scales node-pool + Cloud SQL per tenant.",
           9, False, DARK)],
         [("Indicative GCP list price (Jul 2026); excludes negotiated discounts, taxes & Google program funding · AI lines usage-based.",
           8.5, False, MIDGREY, True)]], line_spacing=1.05, space_after=2)
footer(s, npg())

# ---------- 29. ENGAGEMENT & COMMERCIALS ----------
s = slide()
page_header(s, "Section 5", "Engagement & Next Steps")
textbox(s, Inches(0.55), Inches(1.6), Inches(4.5), Inches(0.3), [[("NEXT STEPS", 12, True, POPPY)]])
bullets(s, Inches(0.55), Inches(1.98), Inches(4.5), Inches(2.9),
        [("One tailored proposal ", "— Devoteam compiles & tailors all your requirements (target: next week)"),
         ("Parallel Google track ", "— Caroline opens the solution discussion with Google"),
         ("You share ", "— report / visualization examples + the Escalation Policy SOP"),
         ("You confirm ", "— proposal meets expectations, then we schedule kick-off")],
        size=11, gap=11)
rect(s, Inches(5.15), Inches(1.62), Inches(4.3), Inches(1.75), AQUA)
textbox(s, Inches(5.35), Inches(1.74), Inches(3.9), Inches(0.3), [[("GOOGLE PROGRAM FUNDING", 11.5, True, DARK)]])
bullets(s, Inches(5.35), Inches(2.14), Inches(3.9), Inches(1.15),
        [("Offsets cost ", "— pursue Google funding to cover development / deployment"),
         ("Leverages ", "your Google Cloud relationship — presented in parallel")],
        size=10.5, gap=8, mcolor=POPPY)
rect(s, Inches(5.15), Inches(3.48), Inches(4.3), Inches(1.4), POPPYLT)
textbox(s, Inches(5.35), Inches(3.6), Inches(3.9), Inches(0.3), [[("KNOWN DELTAS BEING CLOSED", 11.5, True, POPPY)]])
bullets(s, Inches(5.35), Inches(4.0), Inches(3.9), Inches(0.85),
        [("Email · Voice-note · Image ", "— the three items still to complete"),
         ("Everything else ", "aligns to your expectations today")],
        size=10.5, gap=8, mcolor=POPPY)
footer(s, npg())

# ---------- 28. WHY THIS WORKS (dark closing) ----------
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, 0, Inches(0.24), SH, POPPY)
rect(s, 0, Inches(4.86), SW, Inches(0.10), FIRE)
textbox(s, Inches(0.7), Inches(0.72), Inches(8.6), Inches(0.9), [[("Why This Works", 34, True, WHITE)]])
rect(s, Inches(0.72), Inches(1.48), Inches(2.4), Pt(3), FIRE)
bullets(s, Inches(0.75), Inches(1.82), Inches(8.5), Inches(2.6),
        [("Proven live — ", "demonstrated against your SOP on 28 Jul 2026: migration + hardening, not invention"),
         ("Open + owned — ", "no per-seat lock-in; PROTON controls the data, the AI, and the roadmap"),
         ("AI-native on Google — ", "Gemini + Vertex AI, grounded on your knowledge; measurable deflection & faster resolution"),
         ("Cost-down, capability-up — ", "managed infrastructure you scale on your terms, with Google funding in play")],
        size=13, gap=11, color=WHITE, mcolor=FIRE)
textbox(s, Inches(0.7), Inches(4.5), Inches(8.6), Inches(0.32),
        [[("Cost  ·  Ownership  ·  Speed  ·  Proven", 14, True, PINK)]])
logo(s, dark=False, w=Inches(1.5), x=SW - Inches(2.0), y=SH - Inches(0.9))
npg()

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
