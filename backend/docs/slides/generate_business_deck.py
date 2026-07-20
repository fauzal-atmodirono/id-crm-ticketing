"""Generate docs/slides/proton-conversational-ai-business-overview.pptx.

Business/stakeholder companion to the technical deck — benefit-led, no
commands or configuration. Cost figures come from
docs/proposals/2026-06-25-omnichannel-poc-cost-breakdown.md (public rate
cards as of 2026-06-25 — confirm before contractual quoting).

Run from the repo root:
    uv run --with python-pptx python docs/slides/generate_business_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "proton-conversational-ai-business-overview.pptx"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

INK = RGBColor(0x1F, 0x23, 0x28)
MUTED = RGBColor(0x5F, 0x63, 0x68)
ACCENT = RGBColor(0x0B, 0x57, 0xD0)
ACCENT_SOFT = RGBColor(0xE8, 0xF0, 0xFE)
TEAL = RGBColor(0x0F, 0x7B, 0x6C)
TEAL_SOFT = RGBColor(0xE0, 0xF2, 0xEF)
CARD_BORDER = RGBColor(0xDA, 0xDC, 0xE0)
CARD_BG = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARN = RGBColor(0xB0, 0x60, 0x00)
WARN_SOFT = RGBColor(0xFE, 0xF7, 0xE0)

BODY_FONT = "Helvetica Neue"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def _style_run(run, size, color=INK, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = BODY_FONT


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def footer(slide, num):
    tf = textbox(slide, 11.6, 7.08, 1.5, 0.35)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(num)
    _style_run(r, 10, MUTED)


def header(slide, kicker, title):
    tf = textbox(slide, 0.6, 0.32, 12.1, 0.35)
    r = tf.paragraphs[0].add_run()
    r.text = kicker.upper()
    _style_run(r, 12, ACCENT, bold=True)
    tf = textbox(slide, 0.6, 0.62, 12.1, 0.75)
    r = tf.paragraphs[0].add_run()
    r.text = title
    _style_run(r, 27, INK, bold=True)
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.38), Inches(0.55), Pt(3)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False


def bullets(slide, items, x=0.6, y=1.7, w=12.1, size=15.5, gap=9):
    tf = textbox(slide, x, y, w, 7.2 - y)
    first = True
    for item in items:
        text, level, prefix = item if isinstance(item, tuple) else (item, 0, None)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        marker = p.add_run()
        marker.text = "•  " if level == 0 else "–  "
        _style_run(marker, size, ACCENT if level == 0 else MUTED, bold=True)
        if prefix:
            r = p.add_run()
            r.text = prefix + " — "
            _style_run(r, size, INK, bold=True)
        r = p.add_run()
        r.text = text
        _style_run(r, size, INK)
    return tf


def note(slide, text, x=0.6, y=6.35, w=12.1, color=ACCENT, bg=ACCENT_SOFT, label="Note", h=0.62):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.adjustments[0] = 0.12
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"{label}: "
    _style_run(r, 12.5, color, bold=True)
    r = p.add_run()
    r.text = text
    _style_run(r, 12.5, INK)


def table(slide, headers, rows, x, y, w, col_widths=None, size=13, row_h=0.42):
    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(x), Inches(y),
        Inches(w), Inches(row_h * (len(rows) + 1)),
    )
    t = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            t.columns[i].width = Inches(cw)
    for c, htext in enumerate(headers):
        cell = t.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = htext
        _style_run(r, size, WHITE, bold=True)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else CARD_BG
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            _style_run(r, size, INK)
    return t


def card(slide, x, y, w, h, title, body, fill=CARD_BG, border=CARD_BORDER,
         title_color=INK, title_size=14, body_size=11.5, center=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.adjustments[0] = 0.08
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    _style_run(r, title_size, title_color, bold=True)
    if body:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = body
        _style_run(r2, body_size, MUTED if fill in (CARD_BG, WHITE, ACCENT_SOFT, TEAL_SOFT) else WHITE)
    return shp


def flow_arrow(slide, x1, y, x2):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y)
    )
    conn.line.color.rgb = MUTED
    conn.line.width = Pt(1.6)
    return conn


# 1 — Title
s = new_slide()
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()
bar.shadow.inherit = False
tf = textbox(s, 0.9, 2.3, 11.5, 1.8)
r = tf.paragraphs[0].add_run()
r.text = "AI customer service,"
_style_run(r, 42, INK, bold=True)
p = tf.add_paragraph()
r = p.add_run()
r.text = "everywhere your customers are."
_style_run(r, 42, ACCENT, bold=True)
tf = textbox(s, 0.9, 4.35, 11.5, 0.9)
r = tf.paragraphs[0].add_run()
r.text = ("One AI agent that answers on web chat, voice, WhatsApp, phone calls and email — "
          "connected to your knowledge base and your Zendesk.")
_style_run(r, 17, MUTED)
tf = textbox(s, 0.9, 6.7, 11.5, 0.4)
r = tf.paragraphs[0].add_run()
r.text = "Business overview · July 2026"
_style_run(r, 12, MUTED)

# 2 — The opportunity
s = new_slide()
header(s, "Why", "Customers moved. Support tooling didn't.")
bullets(s, [
    ("Customers expect an instant answer at 11 pm on WhatsApp — not a callback window.", 0, "24/7 expectations"),
    ("Most support volume is the same product questions, answered one at a time by people.", 0, "Repetitive load"),
    ("Chat, phone, email and WhatsApp usually mean separate tools, separate queues, separate histories.", 0, "Fragmented channels"),
    ("An unanswered message is not just a support miss — pre-sales questions are where deals quietly die.", 0, "Lost revenue"),
], size=16, gap=14)
note(s, "The goal is not replacing your team — it is letting AI absorb the repetitive 80% "
        "so your people handle the conversations that actually need them.", y=6.2, h=0.75)
footer(s, 2)

# 3 — What we built
s = new_slide()
header(s, "What", "One AI agent, five doors in")
channels = [
    ("Web chat", "on your website"),
    ("Web voice", "tap to talk, spoken replies"),
    ("WhatsApp", "the channel customers prefer"),
    ("Phone", "a real call, answered by AI"),
    ("Email", "replies to your support inbox"),
]
for i, (t, b) in enumerate(channels):
    card(s, 0.6 + i * 2.48, 1.85, 2.28, 1.15, t, b, fill=ACCENT_SOFT, border=ACCENT, title_size=15)
    conn = s.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.6 + i * 2.48 + 1.14), Inches(3.0), Inches(6.67), Inches(3.7),
    )
    conn.line.color.rgb = MUTED
    conn.line.width = Pt(1.4)
card(s, 3.17, 3.7, 7.0, 1.35, "One AI agent",
     "your product knowledge base · your brand voice · your Zendesk",
     fill=ACCENT, border=None, title_color=WHITE, title_size=17, body_size=13)
bullets(s, [
    ("Same brain everywhere: an answer improved once is improved on every channel.", 0, None),
    ("No new tools for your team — everything lands in Zendesk, where they already work.", 0, None),
], y=5.45, size=15)
footer(s, 3)

# 4 — Customer journey
s = new_slide()
header(s, "Experience", "A customer's journey (WhatsApp example)")
steps = [
    ("1 · Asks", "“Does the S70 have adaptive cruise?” — sent at 10:40 pm"),
    ("2 · Answered", "AI replies in seconds with the real spec, from your knowledge base"),
    ("3 · Escalates", "“I want to change my booking” — AI recognises this needs a person"),
    ("4 · Human takes over", "Ticket appears in Zendesk; the agent replies from there — the customer never leaves WhatsApp"),
    ("5 · Handback + CSAT", "Ticket solved → AI resumes and asks for a satisfaction rating"),
]
for i, (t, b) in enumerate(steps):
    x = 0.5 + i * 2.58
    card(s, x, 2.2, 2.38, 2.5, t, b, fill=WHITE, border=ACCENT if i in (0, 1, 4) else WARN,
         title_color=ACCENT if i in (0, 1, 4) else WARN, title_size=14, body_size=11)
    if i < 4:
        flow_arrow(s, x + 2.38, 3.45, x + 2.58)
bullets(s, [
    ("The AI handles steps 1–2 alone, around the clock. Your team only appears at step 4 — and the whole exchange is one continuous conversation for the customer.", 0, None),
], y=5.3, size=15)
footer(s, 4)

# 5 — Web chat + voice
s = new_slide()
header(s, "Channel tour · 1 of 4", "Web chat & voice — your website answers back")
bullets(s, [
    ("Product questions answered instantly, with detailed specs and a visual product carousel.", 0, "Rich answers"),
    ("Customers can tap to talk and hear a natural spoken reply — no typing on mobile.", 0, "Voice built in"),
    ("Every visitor gets a knowledgeable “salesperson”, not a search box.", 0, "Pre-sales impact"),
    ("Fits your site as a standard web component; branding and tone are configurable.", 0, "Your brand"),
], size=15.5, gap=12)
footer(s, 5)

# 6 — WhatsApp
s = new_slide()
header(s, "Channel tour · 2 of 4", "WhatsApp — meet customers in their favourite app")
bullets(s, [
    ("Customers message your business number and get a real answer in seconds — day or night.", 0, "Always on"),
    ("Long answers, spec sheets, follow-up questions — a genuine conversation, not canned menus.", 0, "Not a menu bot"),
    ("Because the AI only replies to inbound messages, Meta's fee for these conversations is zero — you pay fractions of a cent per message.", 0, "Cheap to run"),
    ("Handoff works here too: an agent can take over from Zendesk mid-chat, invisibly to the customer.", 0, "Human backup"),
], size=15.5, gap=12)
footer(s, 6)

# 7 — Phone
s = new_slide()
header(s, "Channel tour · 3 of 4", "Phone — a real call, answered by AI in real time")
bullets(s, [
    ("Customers dial a normal phone number and talk naturally — the AI listens and answers with a human-like voice, live.", 0, "No IVR maze"),
    ("No “press 1 for sales” — callers just say what they need.", 0, "Zero menus"),
    ("Answers come from the same knowledge base as every other channel, so phone answers match what your website says.", 0, "Consistent"),
    ("Multi-language: the AI detects the caller's language automatically (e.g. Malay or English).", 0, "Speaks their language"),
], size=15.5, gap=12)
footer(s, 7)

# 8 — Email
s = new_slide()
header(s, "Channel tour · 4 of 4", "Email — your support inbox, answered")
bullets(s, [
    ("Customers email your support address as usual; the AI drafts and sends the reply through Zendesk.", 0, "Nothing changes for customers"),
    ("Replies arrive in minutes instead of the next business day.", 0, "Speed"),
    ("Optional draft-assist mode: the AI prepares the answer as a private note and an agent approves before it is sent — ideal while building trust.", 0, "Control dial"),
    ("Runs entirely on your existing Zendesk email — no new mail infrastructure.", 0, "Zero new infra"),
], size=15.5, gap=12)
footer(s, 8)

# 9 — Team stays in control
s = new_slide()
header(s, "Control", "Your team stays in control — inside Zendesk")
bullets(s, [
    ("Every AI conversation, on every channel, is logged as a Zendesk ticket your team can read at any time. Nothing is a black box.", 0, "Full visibility"),
    ("The AI recognises frustration, complaints and requests it shouldn't handle — and creates a handoff ticket only then, so agents see real work, not noise.", 0, "Smart escalation"),
    ("Agents reply from Zendesk exactly as they do today; the reply reaches the customer on whichever channel they used.", 0, "No new tools"),
    ("When the agent solves the ticket, the AI takes back over and collects a satisfaction rating.", 0, "Clean handback"),
], size=15.5, gap=12)
footer(s, 9)

# 10 — What you can measure
s = new_slide()
header(s, "Insight", "If it matters, it's measured")
metrics = [
    ("Volume", "conversations per channel, per day"),
    ("Containment", "how much the AI resolves alone"),
    ("Speed", "response time per turn"),
    ("Fallback rate", "how often the AI couldn't answer"),
    ("Bounce", "customers who leave after one turn"),
    ("CSAT", "satisfaction after human handoffs"),
    ("NPS", "loyalty score collected in-chat"),
    ("Quality", "human-reviewed answer accuracy"),
]
for i, (t, b) in enumerate(metrics):
    x = 0.6 + (i % 4) * 3.12
    y = 1.85 + (i // 4) * 1.6
    card(s, x, y, 2.92, 1.4, t, b, title_color=ACCENT, title_size=15, body_size=11.5)
bullets(s, [
    ("All of it lives in a management dashboard (with the raw data in your warehouse for BI tools) — so the AI's business case is visible week by week, not anecdotal.", 0, None),
], y=5.35, size=15)
footer(s, 10)

# 11 — Running costs
s = new_slide()
header(s, "Economics", "What it costs to run")
table(s, ["Channel", "Running cost", "Why"], [
    ["WhatsApp", "≈ $0.005 per message", "Only Twilio's fee — Meta's fee is $0 for replies within the service window"],
    ["Email", "≈ $0 per email", "Rides on your existing Zendesk email; only AI inference"],
    ["Web chat & voice", "≈ $0 per conversation", "AI inference only — fractions of a cent"],
    ["Phone", "≈ $0.03–0.08 per minute", "Call minutes + real-time speech processing (option-dependent)"],
    ["Fixed", "≈ $4 per month", "One local phone number rental (Malaysia)"],
], x=0.6, y=1.75, w=12.1, col_widths=[2.5, 3.0, 6.6], size=13.5, row_h=0.56)
bullets(s, [
    ("A thousand WhatsApp support conversations cost dollars, not headcount. The dominant platform costs are the usual Google Cloud and Zendesk subscriptions you'd run anyway.", 0, None),
], y=4.95, size=14.5)
note(s, "Figures from public Twilio/Meta/Google rate cards as of 25 Jun 2026 (see the cost "
        "breakdown document). Confirm current rates before contractual quoting.",
     y=6.15, label="Disclaimer", color=WARN, bg=WARN_SOFT, h=0.75)
footer(s, 11)

# 12 — Live today
s = new_slide()
header(s, "Status", "This is not a concept — it's running today")
table(s, ["Channel", "Status in the demo environment"], [
    ["Web chat", "Live — real product answers from the knowledge base"],
    ["Web voice", "Live — spoken replies"],
    ["WhatsApp", "Live — message the demo number and the AI answers"],
    ["Phone", "Live — call the demo number and talk to the AI"],
    ["Email", "Live — email the demo address, get an AI reply"],
    ["Dashboard", "Live — metrics visible in the management dashboard"],
], x=0.6, y=1.75, w=12.1, col_widths=[2.6, 9.5], size=13.5, row_h=0.5)
bullets(s, [
    ("Every channel above has been verified end-to-end with real messages, real calls and real emails on the cloud deployment — not a staged video.", 0, None),
], y=5.15, size=14.5)
footer(s, 12)

# 13 — Path to production
s = new_slide()
header(s, "Roadmap", "From demo to your production launch")
steps = [
    ("1 · Brand it", "Your knowledge base onboarded; tone and languages tuned to your brand"),
    ("2 · Own the numbers", "Meta-verified WhatsApp sender and a phone number on your own brand"),
    ("3 · Wire your Zendesk", "Your (or a shared) Zendesk instance, your agents, your workflows"),
    ("4 · Launch & measure", "Go live channel by channel; watch containment and CSAT from week one"),
]
for i, (t, b) in enumerate(steps):
    card(s, 0.6 + i * 3.12, 2.0, 2.92, 2.2, t, b, fill=WHITE, border=ACCENT,
         title_color=ACCENT, title_size=14.5, body_size=11.5)
    if i < 3:
        flow_arrow(s, 0.6 + i * 3.12 + 2.92, 3.1, 0.6 + (i + 1) * 3.12)
bullets(s, [
    ("The demo already exercises every moving part — production is configuration and branding, not a rebuild. A phased launch (e.g. WhatsApp first) de-risks adoption.", 0, None),
], y=4.7, size=15)
footer(s, 13)

# 14 — Live demo flow
s = new_slide()
header(s, "Demo", "What you'll see in the live demo")
bullets(s, [
    ("Ask the website a detailed product question — watch it answer with specs and a product carousel.", 0, "1 · Web chat"),
    ("Ask the same thing by voice and hear the reply.", 0, "2 · Web voice"),
    ("Message the demo WhatsApp number from your own phone.", 0, "3 · WhatsApp"),
    ("Call the demo number and have a spoken conversation with the AI.", 0, "4 · Phone"),
    ("Type an angry message — watch a ticket appear in Zendesk, reply as the agent, and see the customer get it on their channel.", 0, "5 · Handoff"),
    ("Solve the ticket — the AI resumes and asks for a rating.", 0, "6 · Handback + CSAT"),
    ("Open the metrics dashboard and see today's demo traffic in the numbers.", 0, "7 · Dashboard"),
], size=14.5, gap=8)
footer(s, 14)

# 15 — Next steps
s = new_slide()
header(s, "Next steps", "Where we go from here")
bullets(s, [
    ("Pick the launch channels and the first brand/market.", 0, "Decide"),
    ("Share your product content — we onboard the knowledge base and tune the brand voice.", 0, "Prepare"),
    ("A pilot deployment on your accounts is roughly a one-day setup, then a measured pilot period with weekly numbers.", 0, "Pilot"),
    ("Review containment, CSAT and cost per conversation together — then scale to the remaining channels.", 0, "Scale"),
], size=16, gap=14)
tf = textbox(s, 0.6, 6.3, 12.1, 0.6)
r = tf.paragraphs[0].add_run()
r.text = "One AI agent. Five channels. Your team in control — and the numbers to prove it."
_style_run(r, 16, ACCENT, bold=True)
footer(s, 15)

prs.save(OUT)

# ---- verify ----
check = Presentation(OUT)
n = len(check.slides._sldIdLst)  # noqa: SLF001
assert n == 15, f"expected 15 slides, got {n}"
for idx, slide in enumerate(check.slides, 1):
    texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
    assert any(t.strip() for t in texts), f"slide {idx} has no text"
print(f"OK: wrote {OUT} ({n} slides)")
