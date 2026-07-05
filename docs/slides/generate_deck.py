"""Generate docs/slides/proton-cloud-run-deployment-guide.pptx.

Run from the repo root:
    uv run --with python-pptx python docs/slides/generate_deck.py

Part A is a standalone executive summary; Part B is the engineer runbook
(Cloud Run, Twilio, Zendesk, Metrics, Troubleshooting). All identifiers come
from the repo (.env.example, scripts/) and the recorded 2026-06-30..07-02
production deploy; secret VALUES are never included.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "proton-cloud-run-deployment-guide.pptx"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

INK = RGBColor(0x1F, 0x23, 0x28)
MUTED = RGBColor(0x5F, 0x63, 0x68)
ACCENT = RGBColor(0x0B, 0x57, 0xD0)
ACCENT_SOFT = RGBColor(0xE8, 0xF0, 0xFE)
TEAL = RGBColor(0x0F, 0x7B, 0x6C)
TEAL_SOFT = RGBColor(0xE0, 0xF2, 0xEF)
CODE_BG = RGBColor(0xF1, 0xF3, 0xF4)
CARD_BORDER = RGBColor(0xDA, 0xDC, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARN = RGBColor(0xB0, 0x60, 0x00)
WARN_SOFT = RGBColor(0xFE, 0xF7, 0xE0)

BODY_FONT = "Helvetica Neue"
CODE_FONT = "Menlo"

BACKEND_URL = "https://proton-backend-247165654737.asia-southeast1.run.app"
FRONTEND_URL = "https://proton-frontend-247165654737.asia-southeast1.run.app"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def _style_run(run, size, color=INK, bold=False, font=BODY_FONT, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    run.font.italic = italic


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def footer(slide, num):
    tf = textbox(slide, 11.6, 7.08, 1.5, 0.35)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _style_run(p.add_run(), 10, MUTED)
    p.runs[0].text = str(num)


def header(slide, kicker, title):
    tf = textbox(slide, 0.6, 0.32, 12.1, 0.35)
    r = tf.paragraphs[0].add_run()
    r.text = kicker.upper()
    _style_run(r, 12, ACCENT, bold=True)
    tf = textbox(slide, 0.6, 0.62, 12.1, 0.75)
    r = tf.paragraphs[0].add_run()
    r.text = title
    _style_run(r, 27, INK, bold=True)
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.38), Inches(0.55), Pt(3)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False


def bullets(slide, items, x=0.6, y=1.65, w=12.1, size=15):
    """items: list of str or (text, level) or (text, level, bold_prefix)."""
    tf = textbox(slide, x, y, w, 7.2 - y)
    first = True
    for item in items:
        text, level, prefix = item if isinstance(item, tuple) else (item, 0, None)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(7)
        marker = p.add_run()
        marker.text = ("•  " if level == 0 else "–  ")
        _style_run(marker, size, ACCENT if level == 0 else MUTED, bold=True)
        if prefix:
            r = p.add_run()
            r.text = prefix + " — "
            _style_run(r, size, INK, bold=True)
        r = p.add_run()
        r.text = text
        _style_run(r, size, INK)
    return tf


def code_block(slide, lines, x=0.6, y=None, w=12.1, size=12, title=None):
    h = 0.24 * len(lines) + 0.3 + (0.28 if title else 0)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.adjustments[0] = 0.06
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = CARD_BORDER
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_top = Inches(0.1)
    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        first = False
        r = p.add_run()
        r.text = title
        _style_run(r, 11, MUTED, bold=True)
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        _style_run(r, size, INK, font=CODE_FONT)
    return y + h


def note(slide, text, x=0.6, y=6.35, w=12.1, color=WARN, bg=WARN_SOFT, label="Gotcha"):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.62)
    )
    box.adjustments[0] = 0.12
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"{label}: "
    _style_run(r, 12.5, color, bold=True)
    r = p.add_run()
    r.text = text
    _style_run(r, 12.5, INK)


def table(slide, headers, rows, x, y, w, col_widths=None, size=12.5, row_h=0.34):
    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(x), Inches(y),
        Inches(w), Inches(row_h * (len(rows) + 1)),
    )
    t = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            t.columns[i].width = Inches(cw)
    for c, htext in enumerate(headers):
        cell = t.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = htext
        _style_run(r, size, WHITE, bold=True)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xF8, 0xF9, 0xFA)
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            mono = str(val).startswith("`")
            r.text = str(val).replace("`", "")
            _style_run(r, size - (1 if mono else 0), INK, font=CODE_FONT if mono else BODY_FONT)
    return t


def box_shape(slide, x, y, w, h, text, fill, border, text_color, size=12.5, bold=True, sub=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.adjustments[0] = 0.1
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _style_run(r, size, text_color, bold=bold)
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        _style_run(r2, size - 2.5, text_color if fill != WHITE else MUTED)
    return shp


def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = MUTED
    conn.line.width = Pt(1.4)
    return conn


def divider(kicker, title, items=None):
    slide = new_slide()
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT
    bg.line.fill.background()
    bg.shadow.inherit = False
    tf = textbox(slide, 0.9, 2.4, 11.5, 0.5)
    r = tf.paragraphs[0].add_run()
    r.text = kicker.upper()
    _style_run(r, 15, RGBColor(0xBF, 0xD5, 0xF7), bold=True)
    tf = textbox(slide, 0.9, 2.9, 11.5, 1.2)
    r = tf.paragraphs[0].add_run()
    r.text = title
    _style_run(r, 40, WHITE, bold=True)
    if items:
        tf = textbox(slide, 0.95, 4.3, 11.5, 2.5)
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(6)
            r = p.add_run()
            r.text = it
            _style_run(r, 16, WHITE)
    return slide


# ---------------------------------------------------------------- Part A ----

# 1 — Title
s = new_slide()
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()
bar.shadow.inherit = False
tf = textbox(s, 0.9, 2.35, 11.5, 1.8)
r = tf.paragraphs[0].add_run()
r.text = "Proton Conversational AI"
_style_run(r, 44, INK, bold=True)
p = tf.add_paragraph()
r = p.add_run()
r.text = "Running on Cloud Run — deployment & channel integration guide"
_style_run(r, 22, MUTED)
tf = textbox(s, 0.9, 4.6, 11.5, 1.2)
r = tf.paragraphs[0].add_run()
r.text = "Part A · Executive summary (slides 2–6)"
_style_run(r, 15, INK)
p = tf.add_paragraph()
r = p.add_run()
r.text = "Part B · Engineer runbook: Cloud Run, Twilio, Zendesk, Metrics (slides 7–33)"
_style_run(r, 15, INK)
tf = textbox(s, 0.9, 6.7, 11.5, 0.4)
r = tf.paragraphs[0].add_run()
r.text = "July 2026 · repo: github.com/Yudaadi-devo/proton-conversational-ai"
_style_run(r, 12, MUTED)

# 2 — What it does
s = new_slide()
header(s, "Part A · Overview", "One AI agent, five customer channels")
bullets(s, [
    ("Web chat + web voice", 0, "Browser"),
    ("Vue SPA on Cloud Run — text chat with KB answers, product carousel; tap-to-talk voice with Gemini TTS replies.", 1, None),
    ("WhatsApp", 0, "Twilio"),
    ("Customers message the business number; the AI replies on WhatsApp via Twilio.", 1, None),
    ("Phone (voice calls)", 0, "Twilio"),
    ("Real PSTN calls: Twilio Media Streams bridged to Gemini Live for real-time speech-to-speech.", 1, None),
    ("Email", 0, "Zendesk"),
    ("Inbound support emails answered by the AI through Zendesk-native email.", 1, None),
    ("In every channel", 0, "Common core"),
    ("Gemini (Google ADK) answers from the product knowledge base; every conversation is mirrored to Zendesk tickets; the AI detects when a human is needed and hands off, with CSAT on handback; bot metrics stream to BigQuery.", 1, None),
], size=14.5)
footer(s, 2)

# 3 — Architecture
s = new_slide()
header(s, "Part A · Overview", "Architecture")
chan_labels = [
    ("Browser", "chat + voice (SPA)"),
    ("WhatsApp", "via Twilio sandbox/sender"),
    ("Phone (PSTN)", "Twilio Media Streams"),
    ("Email", "Zendesk-native"),
]
def _clamp(v, lo=2.75, hi=4.05):
    return min(max(v, lo), hi)


for i, (label, sub) in enumerate(chan_labels):
    y = 1.85 + i * 1.02
    box_shape(s, 0.6, y, 2.7, 0.8, label, ACCENT_SOFT, ACCENT, INK, sub=sub)
    arrow(s, 3.3, y + 0.4, 4.7, _clamp(y + 0.4))
box_shape(
    s, 4.7, 2.55, 3.5, 1.7,
    "Cloud Run · proton-backend", ACCENT, None, WHITE, size=14,
    sub="FastAPI + Google ADK · webhooks · wss bridge",
)
box_shape(
    s, 4.7, 5.05, 3.5, 0.9,
    "Cloud Run · proton-frontend", TEAL_SOFT, TEAL, INK,
    sub="Vue SPA served by nginx",
)
arrow(s, 6.45, 5.05, 6.45, 4.25)
svc_labels = [
    ("Gemini (Vertex AI)", "chat, TTS, Live audio"),
    ("Vertex AI Search", "product knowledge base"),
    ("Zendesk", "tickets · handoff · email"),
    ("BigQuery", "bot metrics + dashboard"),
]
for i, (label, sub) in enumerate(svc_labels):
    y = 1.85 + i * 1.02
    box_shape(s, 9.9, y, 2.9, 0.8, label, RGBColor(0xF8, 0xF9, 0xFA), CARD_BORDER, INK, sub=sub)
    arrow(s, 8.2, _clamp(y + 0.4), 9.9, y + 0.4)
tf = textbox(s, 0.6, 6.55, 12.1, 0.5)
r = tf.paragraphs[0].add_run()
r.text = ("Everything server-side is one FastAPI service — channel webhooks, the ADK agent, "
          "Zendesk mirroring and the phone wss bridge — so one deploy updates every channel.")
_style_run(r, 12.5, MUTED, italic=True)
footer(s, 3)

# 4 — What going live requires
s = new_slide()
header(s, "Part A · Overview", "Going live needs three vendor accounts")
cards = [
    ("Google Cloud", ACCENT, ACCENT_SOFT, [
        "Cloud Run (backend + frontend)",
        "Vertex AI: Gemini, TTS, Live, AI Search KB",
        "Secret Manager for Twilio secrets",
        "BigQuery for bot metrics",
        "Need: billing-enabled project + gcloud CLI",
    ]),
    ("Twilio", TEAL, TEAL_SOFT, [
        "WhatsApp sender (sandbox for demo, Meta-approved sender for prod)",
        "A voice phone number (inbound calls)",
        "TwiML App + API key for browser softphone",
        "Need: account, auth token, API key pair",
    ]),
    ("Zendesk", WARN, WARN_SOFT, [
        "Tickets: conversation mirror + AI handoff",
        "Email channel (native support address)",
        "Webhooks + triggers back to the backend",
        "Need: Suite account, admin email + API token",
    ]),
]
for i, (name, color, soft, items) in enumerate(cards):
    x = 0.6 + i * 4.15
    box_shape(s, x, 1.75, 3.9, 0.6, name, soft, None, color, size=17)
    tf = textbox(s, x + 0.08, 2.5, 3.8, 4.3)
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        m = p.add_run()
        m.text = "•  "
        _style_run(m, 13.5, color, bold=True)
        r = p.add_run()
        r.text = it
        _style_run(r, 13.5, INK, bold=it.startswith("Need:"))
footer(s, 4)

# 5 — Effort at a glance
s = new_slide()
header(s, "Part A · Overview", "Deployment effort at a glance")
table(s, ["Phase", "What happens", "Rough effort"], [
    ["1 · GCP + deploy", "Enable APIs, deploy backend + frontend to Cloud Run, set env vars + secrets", "1–2 h"],
    ["2 · Twilio", "API keys, WhatsApp sandbox webhook, phone number + TwiML App", "~1 h"],
    ["3 · Zendesk", "API token, relay/CSAT webhooks + triggers, email channel", "~1 h"],
    ["4 · Metrics", "BigQuery dataset + tables, seed or sync, enable per-turn streaming", "~1 h"],
    ["5 · Smoke test", "One live pass per channel: chat, voice, WhatsApp, a phone call, an email", "~1 h"],
], x=0.6, y=1.8, w=12.1, col_widths=[2.4, 7.6, 2.1], size=14, row_h=0.62)
note(s, "About one working day end-to-end; most of it is console wiring, not code. "
        "No CI/CD pipeline exists — deploys are two manual gcloud commands.",
     y=5.7, label="Reality check", color=ACCENT, bg=ACCENT_SOFT)
footer(s, 5)

# 6 — Status today
s = new_slide()
header(s, "Part A · Overview", "Status today (prod, as of 2026-07-02)")
table(s, ["Channel", "Status", "Notes"], [
    ["Web chat", "Live", "KB answers verified on prod"],
    ["Web voice", "Live", "Gemini TTS replies"],
    ["WhatsApp", "Live", "Twilio sandbox number → prod webhook; greetings + KB verified"],
    ["Phone", "Live", "+1 510 896 3929 → Gemini Live bridge; smoke test passed"],
    ["Email", "Live", "Zendesk webhook + trigger → prod backend"],
    ["Metrics", "Live", "BigQuery-backed; in-app /dashboard on prod frontend"],
], x=0.6, y=1.75, w=12.1, col_widths=[2.2, 1.6, 8.3], size=13.5, row_h=0.5)
bullets(s, [
    ("Backend " + BACKEND_URL, 1, None),
    ("Frontend " + FRONTEND_URL, 1, None),
    ("Still open: rotate the reused Twilio API key (hygiene), build Looker tiles, run a real Zendesk→BigQuery sync.", 0, "Pending"),
], y=5.5, size=13.5)
footer(s, 6)

# ---------------------------------------------------------------- Part B ----

divider("Part B", "Engineer runbook", [
    "1 · Cloud Run — deploy backend + frontend, env vars, secrets, redeploys",
    "2 · Twilio — WhatsApp sandbox, phone number, TwiML App, hardening",
    "3 · Zendesk — tickets, webhooks + triggers, email channel",
    "4 · Metrics — BigQuery, dashboards, NPS + QA",
    "5 · Troubleshooting + summary",
])  # 7

# ---- Cloud Run ----
divider("Runbook · Section 1", "Cloud Run")  # 8

s = new_slide()  # 9
header(s, "Runbook · Cloud Run", "Prerequisites")
bullets(s, [
    ("gcloud CLI authenticated against a billing-enabled project (prod uses lv-playground-genai, region asia-southeast1).", 0, "GCP"),
    ("Cloud Run, Vertex AI, Cloud Text-to-Speech, Discovery Engine (Vertex AI Search), Secret Manager, BigQuery.", 0, "Enable APIs"),
    ("apps/backend (Python 3 / FastAPI, own Dockerfile) and apps/frontend (Vue 3 / Vite, multi-stage Dockerfile). Each deploys as its own Cloud Run service straight from source.", 0, "Repo layout"),
    ("The Cloud Run runtime service account needs Vertex AI + BigQuery access and roles/secretmanager.secretAccessor for the Twilio secrets. Prod uses the default compute SA.", 0, "Service account"),
], size=14.5)
code_block(s, [
    "gcloud services enable run.googleapis.com aiplatform.googleapis.com \\",
    "    texttospeech.googleapis.com discoveryengine.googleapis.com \\",
    "    secretmanager.googleapis.com bigquery.googleapis.com",
], y=5.3)
footer(s, 9)

s = new_slide()  # 10
header(s, "Runbook · Cloud Run", "Deploy the backend")
y = code_block(s, [
    "gcloud run deploy proton-backend \\",
    "    --source apps/backend \\",
    "    --region asia-southeast1 --allow-unauthenticated",
], y=1.7, title="From a clean main checkout, repo root")
bullets(s, [
    ("--source builds from apps/backend/Dockerfile via Cloud Build — no registry step needed.", 0, None),
    ("Repeat deploys with --source PRESERVE the service's existing env vars and secrets; only the image changes.", 0, "Important"),
    ("GET / returns {status, crm_provider, voice_provider, model} — the health check to confirm config took.", 0, "Verify"),
], y=y + 0.3, size=14.5)
code_block(s, [
    f"curl {BACKEND_URL}/",
], y=4.6, title="Health check")
note(s, "New revisions may get 0% traffic if traffic was ever pinned. After every deploy run: "
        "gcloud run services update-traffic proton-backend --region asia-southeast1 --to-latest")
footer(s, 10)

s = new_slide()  # 11
header(s, "Runbook · Cloud Run", "Deploy the frontend")
bullets(s, [
    ("The SPA bakes the backend URL at build time from the committed apps/frontend/.env (VITE_API_BASE_URL). Point it at the backend service URL BEFORE deploying.", 0, "Config"),
    ("Multi-stage Dockerfile: node builds the Vite bundle, then nginx:alpine serves it. nginx.conf.template (envsubst at start) includes the SPA fallback so hard-loading /dashboard works.", 0, "Image"),
    ("The backend's FRONTEND_ORIGINS env var must include the frontend URL, or browser calls fail CORS (it also gates the phone-token endpoint).", 0, "CORS"),
], size=14.5)
code_block(s, [
    "# apps/frontend/.env",
    f"VITE_API_BASE_URL={BACKEND_URL}",
    "",
    "gcloud run deploy proton-frontend \\",
    "    --source apps/frontend \\",
    "    --region asia-southeast1 --allow-unauthenticated",
], y=4.55)
footer(s, 11)

s = new_slide()  # 12
header(s, "Runbook · Cloud Run", "Env vars & Secret Manager")
table(s, ["Group", "Key variables (see apps/backend/.env.example for all)"], [
    ["Core", "`CRM_PROVIDER=zendesk · KNOWLEDGE_PROVIDER · FRONTEND_ORIGINS`"],
    ["Gemini", "`GOOGLE_GENAI_USE_VERTEXAI · VERTEX_PROJECT_ID · GEMINI_MODEL · GEMINI_LIVE_MODEL`"],
    ["Zendesk", "`ZENDESK_SUBDOMAIN · ZENDESK_EMAIL · ZENDESK_API_TOKEN · ZENDESK_TICKET_TAG`"],
    ["Twilio", "`TWILIO_ACCOUNT_SID · TWILIO_API_KEY_SID · TWILIO_TWIML_APP_SID · TWILIO_WEBHOOK_BASE_URL`"],
    ["Metrics", "`METRICS_PROVIDER=bigquery · BIGQUERY_PROJECT_ID · BIGQUERY_DATASET`"],
], x=0.6, y=1.7, w=12.1, col_widths=[1.7, 10.4], size=12.5, row_h=0.42)
code_block(s, [
    "printf '%s' \"$TOKEN\" | gcloud secrets create twilio-auth-token --data-file=-",
    "gcloud secrets add-iam-policy-binding twilio-auth-token \\",
    "    --member serviceAccount:<runtime-SA> --role roles/secretmanager.secretAccessor",
    "gcloud run services update proton-backend --region asia-southeast1 \\",
    "    --set-secrets TWILIO_AUTH_TOKEN=twilio-auth-token:latest,\\",
    "TWILIO_API_KEY_SECRET=twilio-api-key-secret:latest",
], y=4.35, title="Secrets: twilio-auth-token + twilio-api-key-secret live in Secret Manager")
note(s, "Zendesk credentials are still plain env vars on the prod service — migrate them to "
        "Secret Manager the same way when convenient.", y=6.45, label="Debt")
footer(s, 12)

s = new_slide()  # 13
header(s, "Runbook · Cloud Run", "Redeploys & current prod values")
code_block(s, [
    "gcloud run deploy proton-backend  --source apps/backend  --region asia-southeast1",
    "gcloud run deploy proton-frontend --source apps/frontend --region asia-southeast1",
    "gcloud run services update-traffic proton-backend --region asia-southeast1 --to-latest",
    "gcloud run services update proton-backend --region asia-southeast1 \\",
    "    --update-env-vars KEY=value       # config-only change, no rebuild",
], y=1.7, title="Routine operations")
table(s, ["Prod value", ""], [
    ["Project / region", "lv-playground-genai (247165654737) · asia-southeast1"],
    ["Backend", BACKEND_URL],
    ["Frontend", FRONTEND_URL],
    ["Runtime SA", "247165654737-compute@developer.gserviceaccount.com"],
    ["Latest known revision", "proton-backend-00033-nql (config switch to new Twilio account)"],
], x=0.6, y=3.75, w=12.1, col_widths=[2.6, 9.5], size=12.5, row_h=0.42)
note(s, "No CI/CD — every deploy is manual. Deploy only from a clean, pushed main so prod matches git.",
     y=6.15, label="Process")
footer(s, 13)

# ---- Twilio ----
divider("Runbook · Section 2", "Twilio — WhatsApp + Phone")  # 14

s = new_slide()  # 15
header(s, "Runbook · Twilio", "Account & credentials")
bullets(s, [
    ("Account SID (AC…) + Auth Token — webhook signature validation and the REST client.", 0, "Console » Account"),
    ("API Key pair (SK… + secret) — signs browser softphone access tokens (JWT).", 0, "Console » API keys"),
    ("TwiML App (AP…) — the Voice URL container the softphone + purchased number point at.", 0, "Console » Voice"),
    ("Auth token and API key secret go to Secret Manager; the SIDs are plain env vars.", 0, "Storage"),
], size=14.5)
table(s, ["Env var", "What it is"], [
    ["`TWILIO_ACCOUNT_SID`", "AC… account identifier"],
    ["`TWILIO_AUTH_TOKEN`", "secret — via Secret Manager"],
    ["`TWILIO_API_KEY_SID` / `TWILIO_API_KEY_SECRET`", "SK… pair for voice tokens (secret half in Secret Manager)"],
    ["`TWILIO_TWIML_APP_SID`", "AP… TwiML App"],
    ["`TWILIO_WEBHOOK_BASE_URL`", "public https base of the backend (prod backend URL)"],
], x=0.6, y=4.35, w=12.1, col_widths=[4.6, 7.5], size=12, row_h=0.42)
footer(s, 15)

s = new_slide()  # 16
header(s, "Runbook · Twilio", "WhatsApp sandbox")
bullets(s, [
    ("Console » Messaging » Try it out » Send a WhatsApp message. Customers join by sending the sandbox code to the shared number +1 415 523 8886.", 0, "Join"),
    ("In Sandbox settings, set “When a message comes in” to POST <backend>/webhooks/twilio-whatsapp. Inbound messages then reach the AI; replies go out via the Twilio REST API.", 0, "Webhook"),
    ("Set TWILIO_WHATSAPP_NUMBER=whatsapp:+14155238886 on the backend.", 0, "Env"),
    ("Production replaces the sandbox with a Meta-approved WhatsApp sender on your own number — same webhook, no code change.", 0, "Prod path"),
], size=14.5)
code_block(s, [
    f"When a message comes in:  POST {BACKEND_URL}/webhooks/twilio-whatsapp",
], y=5.15)
note(s, "The sandbox is SHARED per account and points at ONE URL at a time — pointing it at a local "
        "tunnel for testing silently detaches prod.")
footer(s, 16)

s = new_slide()  # 17
header(s, "Runbook · Twilio", "Phone: number, TwiML App, media stream")
bullets(s, [
    ("Buy a voice-capable number (prod: +1 510 896 3929).", 0, "1"),
    ("Set the TwiML App's Voice URL to POST <backend>/voice/phone/incoming. The backend answers with TwiML: <Connect><Stream url=“wss://<backend>/…”>.", 0, "2"),
    ("Point the number's VoiceApplicationSid at the TwiML App — inbound calls now bridge caller audio ↔ Gemini Live in real time.", 0, "3"),
    ("The in-browser softphone fetches a JWT from POST /voice/phone/token and dials through the same TwiML App.", 0, "4"),
    ("PUBLIC_WSS_BASE_URL only if the wss host differs from TWILIO_WEBHOOK_BASE_URL (it defaults to it, https→wss). GEMINI_LIVE_MODEL=gemini-live-2.5-flash-native-audio — the Vertex publisher id, NOT the AI-Studio “…-preview” name.", 0, "Env"),
], size=14)
note(s, "Set GEMINI_LIVE_LANGUAGE=ms-MY for a Malay-focused demo; unset = auto-detect.",
     y=6.35, label="Tip", color=ACCENT, bg=ACCENT_SOFT)
footer(s, 17)

s = new_slide()  # 18
header(s, "Runbook · Twilio", "Hardening the token endpoint")
bullets(s, [
    ("POST /voice/phone/token is UNAUTHENTICATED by design — a public SPA calls it. Three layers keep it from becoming a free Twilio-minutes faucet:", 0, None),
    ("Origin allowlist — the request Origin must be in FRONTEND_ORIGINS.", 1, None),
    ("Short-lived tokens — PHONE_TOKEN_TTL_SECONDS=300.", 1, None),
    ("Per-IP rate limit — PHONE_TOKEN_RATE_LIMIT=10 per PHONE_TOKEN_RATE_WINDOW_SECONDS=60.", 1, None),
    ("Inbound Twilio webhooks are verified with the X-Twilio-Signature header against TWILIO_AUTH_TOKEN.", 0, "Webhooks"),
], size=14.5)
code_block(s, [
    "curl -X POST -H 'Origin: " + FRONTEND_URL + "' \\",
    f"    {BACKEND_URL}/voice/phone/token   # expect 200 + JWT",
], y=4.9, title="Verify")
footer(s, 18)

s = new_slide()  # 19
header(s, "Runbook · Twilio", "Current prod values")
table(s, ["Item", "Value"], [
    ["Account", "“Demo Proton (new)” · `REDACTED`"],
    ["API key", "`REDACTED` (secret half in Secret Manager)"],
    ["TwiML App", "“Demo Proton” · `AP5c3e6ba4933066d2b765dc94f98ae453` · VoiceUrl → prod /voice/phone/incoming"],
    ["Phone number", "+1 510 896 3929 → VoiceApplicationSid = the TwiML App"],
    ["WhatsApp", "shared sandbox `whatsapp:+14155238886` → prod /webhooks/twilio-whatsapp"],
    ["Webhook base", BACKEND_URL],
], x=0.6, y=1.75, w=12.1, col_widths=[2.4, 9.7], size=12.5, row_h=0.5)
bullets(s, [
    ("Switching the backend to this account intentionally broke the OLD account's WA/phone wiring — the sandbox and secrets now belong to the new account.", 0, "Note"),
    ("The API key was reused from chat during the switch — rotate it for hygiene.", 0, "To do"),
], y=5.1, size=13.5)
footer(s, 19)

# ---- Zendesk ----
divider("Runbook · Section 3", "Zendesk — tickets, webhooks, email")  # 20

s = new_slide()  # 21
header(s, "Runbook · Zendesk", "Access & configuration")
bullets(s, [
    ("Admin Center » Apps and integrations » APIs » add an API token. Auth is email/token basic auth.", 0, "API token"),
    ("Prod instance: devoteam-95614.zendesk.com (shared across demo brands).", 0, "Subdomain"),
    ("Set CRM_PROVIDER=zendesk to activate the adapter.", 0, "Backend"),
], size=14.5)
code_block(s, [
    "CRM_PROVIDER=zendesk",
    "ZENDESK_SUBDOMAIN=devoteam-95614",
    "ZENDESK_EMAIL=<admin email>",
    "ZENDESK_API_TOKEN=<token>            # secret — keep out of git",
    "ZENDESK_REQUESTER_DOMAIN=proton.devoteam.example",
    "ZENDESK_CUSTOMER_NAME_PREFIX=Proton AI Customer",
    "ZENDESK_TICKET_TAG=                   # set per brand on a shared instance",
], y=3.7)
footer(s, 21)

s = new_slide()  # 22
header(s, "Runbook · Zendesk", "Ticket model & shared-instance tenancy")
bullets(s, [
    ("Every conversation is mirrored into a Zendesk ticket as comments; if that ticket gets CLOSED, the adapter detects it (422/404) and rotates to a fresh ticket automatically.", 0, "Conversation mirror"),
    ("The AI creates a handoff ticket only when it detects a human is needed (detection-gated) — agents reply in Zendesk, replies relay to the customer's channel, and solving the ticket hands back to the AI with a CSAT ask.", 0, "Handoff"),
    ("One Zendesk instance can serve multiple brands: per-customer ZENDESK_TICKET_TAG lanes, ZENDESK_REQUESTER_DOMAIN and ZENDESK_CUSTOMER_NAME_PREFIX keep end-users and tickets distinguishable, and triggers route by tag. See docs/zendesk-shared-instance-separation.md.", 0, "Multi-tenancy"),
    ("Ticket identity is per-customer: the Zendesk external_id carries the channel session (e.g. whatsapp-+62…), which the webhook handlers filter on.", 0, "Identity"),
], size=14)
footer(s, 22)

s = new_slide()  # 23
header(s, "Runbook · Zendesk", "Webhooks & triggers (agent relay + CSAT)")
bullets(s, [
    ("Two webhooks point back at the backend, authenticated by an X-Proton-Webhook-Secret header (ZENDESK_SUPPORT_WEBHOOK_SECRET):", 0, None),
    ("agent reply relay → POST /webhooks/zendesk-support (fires on public agent comments)", 1, None),
    ("handback CSAT → POST /webhooks/zendesk-handback (fires on status → solved)", 1, None),
    ("Each webhook is invoked by a trigger. Prefer create-scoped triggers with a tag condition on shared instances — bare triggers are account-wide and fire for every brand's tickets (handlers do filter by external_id, but scoping saves noise).", 0, "Triggers"),
    ("Webhooks/triggers are created via the Zendesk API — see docs/superpowers/specs/2026-06-26-whatsapp-single-ticket-handoff-csat-design.md for payloads.", 0, "Provisioning"),
], size=14)
footer(s, 23)

s = new_slide()  # 24
header(s, "Runbook · Zendesk", "Email channel")
bullets(s, [
    ("Customers email the Zendesk support address; a webhook forwards the ticket to the backend, the AI answers, and the reply goes out as a normal Zendesk email — no separate mail infra.", 0, "Flow"),
    ("scripts/provision_zendesk_email_webhook.py creates the webhook + trigger pair against a target base URL (idempotent — safe to rerun after a URL change).", 0, "Provision"),
    ("EMAIL_DRAFT_ASSIST=true makes the AI post replies as PRIVATE notes for an agent to review and send, instead of auto-emailing.", 0, "Draft mode"),
], size=14.5)
code_block(s, [
    "cd apps/backend",
    "uv run python scripts/provision_zendesk_email_webhook.py \\",
    f"    --base-url {BACKEND_URL}",
], y=4.5)
note(s, "Transient webhook failures are not retried by the backend today — a missed email needs a "
        "manual nudge (known gap from the 2026-06-29 smoke test).")
footer(s, 24)

s = new_slide()  # 25
header(s, "Runbook · Zendesk", "Current prod values")
table(s, ["Object", "ID", "Points at"], [
    ["Webhook · agent reply relay", "`01KW1P1PXSZ95MPX3DSDGA08RS`", "prod /webhooks/zendesk-support"],
    ["Webhook · handback CSAT", "`01KW1P1QN8GVA9S6SACH33D6D6`", "prod /webhooks/zendesk-handback"],
    ["Webhook · email channel", "`01KWADDCYCTSCW0T8RQ8Q3VPK1`", "prod email endpoint"],
    ["Trigger · relay public reply", "`16621232702223`", "relay webhook"],
    ["Trigger · CSAT on solved", "`16621257673615`", "CSAT webhook"],
    ["Trigger · email channel", "`16655564058511`", "email webhook"],
], x=0.6, y=1.75, w=12.1, col_widths=[3.9, 4.1, 4.1], size=12.5, row_h=0.48)
bullets(s, [
    ("All webhook URLs were repointed from tunnel URLs to the prod backend at go-live. If they're ever repointed for local testing, they must be PATCHed back or relay/CSAT/email silently stop.", 0, "Caveat"),
    ("To tear down: delete the triggers, then the webhooks, by the IDs above.", 0, "Teardown"),
], y=4.85, size=13.5)
footer(s, 25)

# ---- Metrics ----
divider("Runbook · Section 4", "Metrics — BigQuery + dashboards")  # 26

s = new_slide()  # 27
header(s, "Runbook · Metrics", "BigQuery provisioning")
table(s, ["Table", "Contents", "Filled by"], [
    ["`conversations`", "one row per conversation (channel, timings, handoff, CSAT)", "Zendesk sync"],
    ["`turn_events`", "one row per turn (latency, fallback, bounce)", "backend streaming"],
    ["`qa_labels`", "manual accuracy/quality labels", "POST /qa/label"],
], x=0.6, y=1.7, w=12.1, col_widths=[2.4, 6.0, 3.7], size=13, row_h=0.46)
code_block(s, [
    "cd apps/backend",
    "uv run python scripts/seed_demo_metrics.py       # demo data (625 convs seeded on prod)",
    "uv run python scripts/sync_zendesk_metrics.py    # real Zendesk -> conversations",
], y=3.75, title="Dataset demo_proton in lv-playground-genai; seed or sync:")
note(s, "The trial Zendesk sandbox blocks some zendesk.com API paths — the real sync has not run "
        "on prod yet; the dashboard currently shows seeded demo data.", y=5.35, label="Known gap")
footer(s, 27)

s = new_slide()  # 28
header(s, "Runbook · Metrics", "Providers & the sync scheduler")
bullets(s, [
    ("METRICS_PROVIDER=bigquery streams one row per turn into turn_events (default noop = off). This powers speed, fallback-rate and bounce metrics.", 0, "Per-turn"),
    ("METRICS_SYNC_ENABLED=true refreshes the conversations table from Zendesk every METRICS_SYNC_INTERVAL_HOURS=6 while the backend runs — same creds as the manual sync.", 0, "Scheduler"),
    ("QA_PROVIDER=bigquery enables the quality-label write path.", 0, "QA"),
    ("Prod runs METRICS_PROVIDER=bigquery with BIGQUERY_PROJECT_ID=lv-playground-genai, BIGQUERY_DATASET=demo_proton.", 0, "Prod"),
], size=14.5)
footer(s, 28)

s = new_slide()  # 29
header(s, "Runbook · Metrics", "Dashboards")
bullets(s, [
    ("GET /metrics/dashboard aggregates all 8 metric blocks from BigQuery; the Vue /dashboard page renders it (live on the prod frontend — volume by channel, CSAT, speed, quality…).", 0, "In-app (live)"),
    ("Looker Studio tiles over the same dataset — designs in docs/dashboards/looker-bot-metrics-phase{1,2}.md; the tiles themselves are still to be built.", 0, "Looker (pending)"),
    ("The v_quality view joins qa_labels for the accuracy/quality block.", 0, "Views"),
], size=14.5)
code_block(s, [
    f"open {FRONTEND_URL}/dashboard",
    f"curl {BACKEND_URL}/metrics/dashboard | jq keys",
], y=4.55)
footer(s, 29)

s = new_slide()  # 30
header(s, "Runbook · Metrics", "NPS & QA labelling")
bullets(s, [
    ("POST /chat/nps records an end-of-chat NPS score from the SPA — unauthenticated, written to BigQuery.", 0, "NPS"),
    ("POST /qa/label records a manual accuracy/quality verdict per conversation. Gated by X-API-Key = QA_API_KEY; an EMPTY QA_API_KEY locks the endpoint entirely.", 0, "QA labels"),
], size=14.5)
code_block(s, [
    f"curl -X POST {BACKEND_URL}/qa/label \\",
    "    -H 'X-API-Key: <QA_API_KEY>' -H 'Content-Type: application/json' \\",
    "    -d '{\"conversation_id\": \"…\", \"accurate\": true, \"quality\": 4}'",
], y=3.85)
footer(s, 30)

# ---- Troubleshooting + summary ----
s = new_slide()  # 31
header(s, "Runbook · Troubleshooting", "Deploy & platform gotchas")
table(s, ["Symptom", "Cause → fix"], [
    ["Deployed but nothing changed", "Traffic pinned to an old revision → `gcloud run services update-traffic <svc> --to-latest`"],
    ["Phone bridge dies at connect", "AI-Studio “…-preview” Live model id — Vertex rejects it → use `gemini-live-2.5-flash-native-audio`"],
    ["Secrets “missing” after deploy", "They aren't — `--source` preserves env + secrets. Check the SA has `secretAccessor`"],
    ["Browser calls fail CORS / token 403", "Frontend URL missing from `FRONTEND_ORIGINS` on the backend"],
], x=0.6, y=1.75, w=12.1, col_widths=[3.7, 8.4], size=13, row_h=0.72)
footer(s, 31)

s = new_slide()  # 32
header(s, "Runbook · Troubleshooting", "Channel gotchas")
table(s, ["Symptom", "Cause → fix"], [
    ["WhatsApp: short replies work, long ones vanish", "Twilio 1600-char limit (error 21617) → backend now chunks replies; if it recurs check Twilio error logs"],
    ["Zendesk mirror stops updating", "Conversation ticket was CLOSED (422) → backend auto-rotates to a fresh ticket since rev 00032; older builds stall"],
    ["WhatsApp inbound dead after local testing", "Shared sandbox points at ONE url — someone aimed it at a tunnel → set it back to the prod webhook"],
    ["Twilio signature validation fails behind a tunnel", "A LEADING SPACE pasted into the webhook URL changes the signed URL → re-paste without whitespace"],
    ["Email replies never arrive", "Transient webhook failure — no retry today → re-trigger by updating the ticket, and watch backend logs"],
], x=0.6, y=1.75, w=12.1, col_widths=[4.6, 7.5], size=12.5, row_h=0.78)
footer(s, 32)

s = new_slide()  # 33
header(s, "Wrap-up", "Go-live checklist & where things live")
bullets(s, [
    ("Deploy backend + frontend → set env vars + secrets → wire Twilio (sandbox webhook, TwiML App, number) → wire Zendesk (token, 3 webhooks + triggers) → enable metrics → smoke-test all five channels.", 0, "Checklist"),
    ("github.com/Yudaadi-devo/proton-conversational-ai — apps/backend, apps/frontend.", 0, "Repo"),
    ("docs/USAGE.md and docs/USER_GUIDE.md (operation), docs/decisions/000{1,2,3}-*.md (stack ADRs), docs/testing/*-smoke-test.md (channel test scripts), apps/backend/.env.example (every env var, commented).", 0, "Docs"),
    ("Rotate the Twilio API key, build the Looker tiles, run a real Zendesk→BigQuery sync, migrate Zendesk creds to Secret Manager.", 0, "Open items"),
], size=14.5)
tf = textbox(s, 0.6, 6.5, 12.1, 0.5)
r = tf.paragraphs[0].add_run()
r.text = "One repo · one deploy per service · five live channels."
_style_run(r, 15, ACCENT, bold=True)
footer(s, 33)

prs.save(OUT)

# ---- verify ----
check = Presentation(OUT)
n = len(check.slides._sldIdLst)  # noqa: SLF001
assert n == 33, f"expected 33 slides, got {n}"
for idx, slide in enumerate(check.slides, 1):
    texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
    assert any(t.strip() for t in texts), f"slide {idx} has no text"
print(f"OK: wrote {OUT} ({n} slides)")
